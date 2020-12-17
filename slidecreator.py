from pptx.util import Inches, Pt
from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import ChartData
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.slide import SlideLayout
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
import logging

# Internal modules
import chartcreator as chc
import variables as v


def longestval(lst):
    longest_val = lst[0]
    shortest_val = lst[0]
    for val in lst:
        if len(val) >= len(longest_val):
            longest_val = val
        elif len(val) <= len(shortest_val):
            shortest_val = val
    return longest_val, shortest_val


def notesinsert(slide, notestext):
    text_frame = slide.notes_slide.notes_text_frame
    text_frame.text = str(notestext)


def scrub_formatting(text):
    scrubbedtext = str(text)
    for text_format in ['/b', '/i', '/h1', '/h2', '/h3', '/h4', '/h5', '/h6', '/h7',
                        '/q', '/tag', '#', '/mint', '/mandarin']:
        scrubbedtext = scrubbedtext.replace(text_format, '', -1)
    return scrubbedtext


def text_formatting(text, text_unit):
    if '/b' in text:
        text_unit.font.bold = True
    if '/tag' in text:
        text_unit.font.color.theme_color = MSO_THEME_COLOR.BACKGROUND_2
        text_unit.font.color.brightness = -0.5
    for level in ['/h1', '/h2', '/h3', '/h4', '/h5', '/h6', '/h7']:
        if level in text:
            text_unit.level = int(level[-1]) - 1
    if '/mandarin' in text:
        text_unit.font.color.theme_color = v.brand_colors['mandarin']
    if '/mint' in text:
        text_unit.font.color.theme_color = v.brand_colors['mint']


def insert_text(textlst, text_frame, one_level=False):
    for idx, para in enumerate(textlst):
        try:
            p = text_frame.paragraphs[idx]
        except IndexError:
            p = text_frame.add_paragraph()

        p.level = 0
        if not one_level:
            if idx > 0 or '/tag' in para:
                p.level = 1

        # Command-based formatting
        if '#' in para:
            # Break para into parts for formatting
            para_runs = para.split('#')
            for r in para_runs:
                run = p.add_run()
                text_formatting(text=r, text_unit=run)
                run.text = scrub_formatting(r)
        else:
            text_formatting(text=para, text_unit=p)
            p.text = scrub_formatting(para)


def preflightaddshape(slide, msg, lvl, boxplacement):
    left = top = Inches(boxplacement)
    width = height = Inches(2.5)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill = shape.fill
    fill.solid()
    if lvl == 'high':
        fill.fore_color.rgb = RGBColor(255, 0, 0)
    else:
        fill.fore_color.rgb = RGBColor(255, 255, 0)
    line = shape.line
    line.fill.background()
    shape.shadow.inherit = False
    shape.text = msg
    shape.text_frame.paragraphs[0].font.color.theme_color = MSO_THEME_COLOR.TEXT_1
    shape.text_frame.paragraphs[0].font.color.size = Pt(10)


def preflight(slide, el):  # Inserts applicable error messages as boxes
    if len(el) > 0:
        for chart in el:
            if len(chart) == 0:
                logging.info("No Errors Found")
            else:
                boxplacement = 1.0
                for error in chart:
                    if error in v.error_dict:
                        preflightaddshape(slide, v.error_dict[error][0], v.error_dict[error][1], boxplacement)
                        logging.warning(v.error_dict[error][0])
                    boxplacement += 0.5


def add_gradient_legend(slide):
    width, height, top = Inches(1.35), Inches(0.33), Inches(0.47)
    legend_text = ['Low Index', 'High Index']
    left_loc = [10.13, 11.47]

    for text, loc in zip(legend_text, left_loc):
        left = Inches(loc)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        text_frame = shape.text_frame
        p = text_frame.paragraphs[0]
        p.text = text
        p.font.name = 'Arial'
        p.font.size = Pt(11)
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        shape.line.fill.background()
        shape.shadow.inherit = False
        fill = shape.fill
        fill.gradient()
        fill.gradient_angle = 0
        gradient_stops = fill.gradient_stops
        if 'Low' in text:
            gradient_stops[0].color.rgb = RGBColor(236, 124, 37)
            gradient_stops[0].position = 0.25
            gradient_stops[1].color.theme_color = MSO_THEME_COLOR.BACKGROUND_1
            p.alignment = PP_ALIGN.LEFT
        else:
            gradient_stops[0].color.theme_color = MSO_THEME_COLOR.BACKGROUND_1
            gradient_stops[1].color.rgb = RGBColor(163, 119, 182)
            gradient_stops[1].position = 0.25
            p.alignment = PP_ALIGN.RIGHT


def create_footer(slide, footercopy):
    footer = None
    pageno = None
    for shape in slide.placeholders:
        if "SLIDE_NUMBER" in str(shape.placeholder_format.type):  # Master must have placeholder to function
            pageno = slide.placeholders[shape.placeholder_format.idx]
        elif "FOOTER" in str(shape.placeholder_format.type):  # Master must have placeholder to function
            footer = slide.placeholders[shape.placeholder_format.idx]

    footer_text_frame = footer.text_frame
    footer_pageno = pageno.text_frame.paragraphs[0]._p
    # edits XML directly—not ideal, but can't be updated until Python-pptx adds support.
    fld_xml = (
        '<a:fld %s id="{1F4E2DE4-8ADA-4D4E-9951-90A1D26586E7}" type="slidenum">\n'
        '  <a:rPr lang="en-US" smtClean="0"/>\n'
        '  <a:t>2</a:t>\n'
        '</a:fld>\n' % nsdecls("a")
    )
    fld = parse_xml(fld_xml)
    footer_pageno.append(fld)
    footer_a = 'Note: Due to small base, data is directional'

    for para_idx, para in enumerate(footercopy):
        new_para = para.replace("'", "", -1)
        footercopy[para_idx] = new_para

    insert_text(footercopy, footer_text_frame, one_level=True)


def footer_patch(self):
    for ph in self.placeholders:
        yield ph


SlideLayout.iter_cloneable_placeholders = footer_patch  # replaces module code with redefined latent placeholder code


def count_checker(original_data, body=0, title=0, chart=0, table=0, picture=0):
    collected_values = [body, title, chart, table, picture]
    checkable_values = ['body count', 'title count', 'chart count', 'table count', 'picture count']

    # Collect statements to see if meets parameters
    checklist = []
    for value_idx, value in enumerate(checkable_values):
        checklist.append(original_data[value] == collected_values[value_idx])

    meets_parameters = set(checklist) == {True}
    return meets_parameters


def assign_layout(page_config, template_data):  # TO UPDATE, CREATE DICTIONARY FOR FORMAT OPTIONS
    layout_selection, shapes_list = None, None

    for layout in template_data:
        config = template_data[layout]['layout config']
        slide_function = page_config['function']
        try:
            fp = v.slide_function_options[slide_function]
        except KeyError:
            fp = None

        table_count = page_config['number of tables']
        chart_count = page_config['number of charts']
        total_count = table_count + chart_count

        if slide_function is not None:
            if count_checker(config, body=fp['body'], title=fp['title'], chart=fp['chart'],
                             table=fp['table'], picture=fp['picture']):
                layout_selection, shapes_list = layout, template_data[layout]

        else:
            if page_config['has stat']:
                if count_checker(config, body=3, table=table_count, chart=(chart_count - 1)):
                    layout_selection, shapes_list = layout, template_data[layout]
            else:
                if total_count <= 4 or total_count % 2 == 0:
                    if count_checker(config, body=2, table=table_count, chart=chart_count):
                        layout_selection, shapes_list = layout, template_data[layout]
                else:
                    if count_checker(config, body=2, table=table_count, chart=(chart_count + 1)):
                        layout_selection, shapes_list = layout, template_data[layout]

    return layout_selection, shapes_list


def assign_chart_data(df, config):
    if 'DataFrame' in str(type(df)):
        chart_data = ChartData()
        chart_data.categories = df.index
        percent_format = '0'
        if config['dec places'] > 0:
            percent_format += '.'
            for place in range((config['dec places'])):
                percent_format += '0'
        percent_format += '%'
        config['number format'] = percent_format

        # Move to Data Cleanup
        if config['*FORCE FLOAT'] or config['*MEAN']:
            config['number format'] = '0.0'
        elif config['*FORCE INT']:
            config['number format'] = '0'
        elif config['*FORCE CURRENCY']:
            config['number format'] = '$0.00'

        for col in df.columns:
            chart_data.add_series(col, df[col], config['number format'])

    else:
        chart_data = df
    return chart_data


def callout_formatting(shape, run_text, alignment=PP_ALIGN.CENTER, text_color='default'):
    text_frame = shape.text_frame
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = run_text
    font = p.font
    font.size = Pt(11)
    font.bold = text_color != 'default'
    font.color.theme_color = v.brand_colors[text_color]
    if text_color == 'mint':
        font.color.brightness = -0.25


def data_import(app, template, wbdata, templatedata, trusave, msg='Processing Data. '):
    logging.info('Starting Data Import')
    global prs
    prs = Presentation(template)

    # Select Layout and drop in charts
    for slideidx, page in enumerate(wbdata):
        page_config = wbdata[page]['page config']
        footer_text = []

        total_charts = str(page_config['number of charts'] + page_config['number of tables'])
        logging.info('Slide no. ' + str(slideidx + 1) + '--There is/are ' + total_charts + ' Charts:')

        layout, shapeslist = assign_layout(page_config, templatedata)
        used_ph, used_data = [], []  # Used for multiple charts. Checks if placholders/data used already
        slide = prs.slides.add_slide(prs.slide_layouts[layout])
        slidenotes, preflightlst = [], []

        try:  # If page doesn't have a title, drops in section tags
            slide.shapes.title.text = page_config['page title']

        except AttributeError:
            for shape in shapeslist:
                if 'BODY' in shape:
                    if 12 > shapeslist[shape]['width'] > 9:
                        if shapeslist[shape]['height'] < 1:
                            phidx = shapeslist[shape]['index']
                            placeholder = slide.placeholders[phidx]
                            text_frame = placeholder.text_frame
                            p = text_frame.paragraphs[0]
                            p.text = page_config['section tag']

        for chart in wbdata[page]:
            if chart != 'page config':
                chart_config = wbdata[page][chart]['config']
                intended_chart = chart_config['intended chart']
                title_question = chart_config['title question']
                error_list = chart_config['error list']
                slidenotes.append(chart_config['notes'])

                try: # If title question has a value, this will handle it appropriately
                    chart_config['data question'], chart_config['chart title'] = longestval(title_question)
                    footer_text.append('Q: ' + (chart_config['data question']))
                except IndexError: # Some imports skip this step
                    pass

                try: # Adds worksheet information to logging if available
                    msg += str(chart_config['notes'][0])[11:-1]
                except IndexError:
                    pass

                v.log_entry(msg, app_holder=app, fieldno=2)

                if len(chart_config['bases']) > 0:
                    footer_text.append('Base: ' + (str(chart_config['bases'])[1:-1]))
                if len(chart_config['note']) > 0:
                    for item in chart_config['note']:
                        footer_text.append(item)
                df = wbdata[page][chart]['frame']
                chart_data = assign_chart_data(df, chart_config)
                preflightlst.append(error_list)

                if intended_chart == 'TABLE':
                    for shape in shapeslist:
                        if 'TABLE' in shape:
                            if shape not in used_ph and chart_data not in used_data:
                                phidx = shapeslist[shape]['index']
                                placeholder = slide.placeholders[phidx]
                                chc.create_table(df, placeholder, chart_config)
                                if chart_config['*HEAT MAP'] is True:
                                    add_gradient_legend(slide)
                                used_ph.append(shape)
                                used_data.append(chart_data)
                elif intended_chart == 'STAT':
                    for shape in shapeslist:
                        if 'BODY' in shape:
                            if shape not in used_ph and chart_data not in used_data:
                                if shapeslist[shape]['height'] < 3 and shapeslist[shape]['width'] < 3:
                                    phidx = shapeslist[shape]['index']
                                    placeholder = slide.placeholders[phidx]
                                    text_frame = placeholder.text_frame
                                    insert_text(chart_data, text_frame)
                                    used_ph.append(shape)
                                    used_data.append(chart_data)
                elif intended_chart == 'PICTURE':
                    for shape in shapeslist:
                        if 'PICTURE' in shape:
                            if shape not in used_ph and chart_data not in used_data:
                                phidx = shapeslist[shape]['index']
                                placeholder = slide.placeholders[phidx]
                                picture = placeholder.insert_picture(chart_data)
                                used_ph.append(shape)
                                used_data.append(chart_data)
                else:
                    for shape in shapeslist:
                        if 'CHART' in shape:
                            if shape not in used_ph:
                                if chart_data not in used_data:
                                    phidx = shapeslist[shape]['index']
                                    placeholder = slide.placeholders[phidx]
                                    chc.create_chart(df, placeholder, chart_data, chart_config)
                                    used_ph.append(shape)
                                    used_data.append(chart_data)

            else:  # If it is page config
                pagecontent = wbdata[page][chart]
                if len(pagecontent) > 0:
                    textlst = wbdata[page][chart]['page copy']
                    for shape in shapeslist:
                        if 'BODY' in shape:
                            if shape not in used_ph and textlst not in used_data:
                                if (shapeslist[shape]['height'] > 3) or (shapeslist[shape]['height'] > 2 and shapeslist[shape]['width'] > 9):
                                    phidx = shapeslist[shape]['index']
                                    placeholder = slide.placeholders[phidx]
                                    text_frame = placeholder.text_frame
                                    insert_text(textlst, text_frame)
                                    used_ph.append(shape)
                                    used_data.append(textlst)
                if len(pagecontent['callouts']) > 0:
                    shapes = slide.shapes
                    for callout in pagecontent['callouts']:
                        c = pagecontent['callouts'][callout]
                        width, height, left, top = Inches(c[0]), Inches(c[1]), Inches(c[2]), Inches(c[3])
                        shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
                        shape.shadow.inherit = False
                        shape.line.fill.background()
                        fill = shape.fill
                        if c[4] is None:
                            fill.background()
                        elif c[4] == 'white':
                            fill.solid()
                            fill.fore_color.theme_color = MSO_THEME_COLOR.BACKGROUND_1
                            fill.fore_color.brightness = 1

                        if c[6] == 'c':
                            alignment = PP_ALIGN.CENTER
                        else:
                            alignment = PP_ALIGN.LEFT
                        callout_formatting(shape, c[5], alignment, c[7])

        try:
            create_footer(slide, footer_text)
        except AttributeError:
            logging.info('No Footer Placeholder found')

        notesinsert(slide, slidenotes)  # Add notes to every slide
        preflight(slide, preflightlst)  # Add boxes with error messages to applicable slides
    prs.save(trusave)
    v.log_entry('File Saved', app_holder=app, fieldno=2)
