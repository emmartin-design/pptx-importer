from pptx import Presentation
import variables as v


def assign_preferred_layouts(template_data, report_type):
    new_template_data = {}
    for layout in template_data:
        config = template_data[layout]['layout config']
        try:  # If the report type is in the static report dict, will assign preferred status
            if config['name'] in v.preferred_layouts[report_type]:
                config['preferred'] = True
        except KeyError:
            if config['body count'] == 2:  # Section tag and body copy always included
                for count in ['chart count', 'table count']:
                    if config[count] == 1:
                        for placeholder in template_data[layout]:
                            if 'CHART' in placeholder or 'TABLE' in placeholder:
                                if template_data[layout][placeholder]['width'] == 9.12:
                                    config['preferred'] = True
                    elif config[count] > 1:
                        # Preferred layouts have placeholders of the same width.
                        if len(set(config['width test'])) == 1:
                            if max(config['width test']) < 9.12:
                                config['preferred'] = True

        # If preferred, will pass on to other functions, otherwise dropped
        if config['preferred']:
            new_template_data[layout] = template_data[layout]

    return new_template_data


def collect_template_data(template, report_type):
    template_config = {}
    prs = Presentation(template)

    # Create dictionary for each layout
    for layout_idx, layout in enumerate(prs.slide_layouts):
        template_config[layout_idx] = {}
        template_config[layout_idx]['layout config'] = v.assign_layout_config(name=layout.name)
        config = template_config[layout_idx]['layout config']

        # Grab placeholder details shape by shape
        for shape in layout.placeholders:
            shape_index = shape.placeholder_format.idx
            shape_type = str(shape.placeholder_format.type).split()[0]

            for ph_type in v.placeholder_types:
                if ph_type in shape_type:
                    shape_name = shape_type + ' ' + str(shape_index)
                    template_config[layout_idx][shape_name] = {
                        'index': shape_index,
                        'width': round(shape.width.inches, 2),
                        'height': round(shape.height.inches, 2),
                        'left': round(shape.left.inches, 2)
                    }

                    try:  # Increases the placeholder counts for each layout if appropriate
                        config[(shape_type.lower() + ' count')] += 1
                        if shape_type in ['CHART', 'TABLE', 'PICTURE']:
                            config['width test'].append(round(shape.width.inches, 2))
                    except KeyError:
                        pass

    template_config = assign_preferred_layouts(template_config, report_type)
    v.log_entry('Template Analyzed')

    return template_config
