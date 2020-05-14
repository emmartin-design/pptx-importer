# -*- coding: utf-8 -*-
from pptx.util import Inches, Pt
from pptx import Presentation
from appJar import gui
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LABEL_POSITION, XL_LEGEND_POSITION, XL_CHART_TYPE, XL_TICK_MARK, XL_DATA_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import ChartData
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.slide import SlideLayout
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.oxml.xmlchemy import OxmlElement
from openpyxl import load_workbook
from openpyxl.descriptors.excel import HexBinary, ExtensionList
from openpyxl.descriptors.serialisable import Serialisable
from openpyxl.styles import PatternFill
from openpyxl.styles.colors import Color
from openpyxl import Workbook
from collections import OrderedDict
from copy import copy
from datetime import datetime
from datetime import date
from datetime import time
import os
import glob
import logging
import webbrowser
import re
import pandas as pd
import numpy as np

# Version variable for future automatic updates

version = '1.0.0 Beta'

# Creates a new log for each day the software is run
# This makes it more accessible, considering the low volume
# If volume increases, may switch to one file, and use fileneame in log.

loggingfilename = 'gen_' + datetime.now().strftime('%y%B%d') + '.log'
log_dir = os.path.join(os.path.normpath(os.getcwd()), 'logs')
logfn = os.path.join(log_dir, loggingfilename)
logging.basicConfig(filename=logfn, level=logging.DEBUG, format='%(lineno)d:%(levelname)s:%(message)s')


# template filename shortens app functions
templatename = '/Users/emartin/Desktop/Projects/Data Import Projects/General Import/Templates/DATAIMPORT.pptx'
prs = Presentation(templatename)

# May be outmoded by new template
brand_colors = [
    RGBColor(0, 132, 200),
    RGBColor(232, 44, 42),
    RGBColor(67, 176, 42),
    RGBColor(255, 205, 0),
    RGBColor(144, 99, 205),
    RGBColor(255, 130, 0),
    RGBColor(0, 164, 153),
    RGBColor(245, 153, 177),
    RGBColor(191, 215, 48),
    RGBColor(180, 85, 160),
    RGBColor(107, 139, 213),
    RGBColor(176, 42, 48),
    RGBColor(88, 208, 156),
    RGBColor(88, 88, 88),  # Placeholder color
    RGBColor(40, 40, 40),  # Placeholder color
    RGBColor(230, 230, 230),  # Placeholder color
    RGBColor(50, 50, 50),  # Placeholder color
    RGBColor(210, 210, 210),  # Placeholder color
    RGBColor(0, 0, 0),  # Placeholder color
    RGBColor(99, 99, 99),  # Placeholder color
    RGBColor(30, 30, 30),  # Placeholder color
    RGBColor(220, 220, 220),  # Placeholder color
    RGBColor(60, 60, 60),  # Placeholder color
    RGBColor(200, 200, 200),  # Placeholder color
    RGBColor(10, 10, 10)
]


charttypelist = [
    'COLUMN',
    'STACKED COLUMN',
    '100% STACKED COLUMN',
    'PIE',
    'BAR',
    'STACKED BAR',
    'OVERALL STACKED BAR',
    '100% STACKED BAR',
    'LINE',
    'TABLE',
    'HEATMAP',
    'TOP2BOX'
]

# 'Consumer KPIs', 'Menu Scorecards' to be added
reporttypes = [
    'General Import',
    'Global Navigator Country Reports'
]

# For country reports only
countrylist = [
    'Argentina',
    'Australia',
    'Brazil',
    'Canada',
    'Chile',
    'China',
    'Colombia',
    'France',
    'Germany',
    'India',
    'Indonesia',
    'Japan',
    'Malaysia',
    'Mexico',
    'Philippines',
    'Russia',
    'Saudi Arabia',
    'Singapore',
    'South Korea',
    'South Africa',
    'Spain',
    'Thailand',
    'United Arab Emirates',
    'United Kingdom',
    'United States'
]

# Separate from logging, places error notes in the deck
errordict = {
    '100': ['ERROR: Data Selection Error. Chart Import Aborted', 'high'],
    '103': ['ERROR: No Data Question selected.', 'med'],
    '104': ['ERROR: No Data Base selected. Please select', 'med'],
    '105': ['ERROR: Data cannot be sorted for data imports of this configuration', 'med'],
    '106': ['NOTE: Blank Cells Selected in original data.', 'med'],
    '107': ['NOTE: Formulas detected in cells. Data not read.', 'high'],
    '404': ['ERROR: Chart type not recognized.', 'high'],
    '405': ['NOTE: Chart type was automatically selected.', 'high'],
    '501': ['ERROR: No data selected in worksheet', 'high'],
    '801': ['ERROR: Chart placeholder type not found. Import aborted.', 'high']
}

# Dummy data to insert when major data errors arise.
data_error = {
    'categories': ['data mismatch', 'error', 'data mismatch', 'error'],
    'error': [0, 0, 0, 0]
}


def general_formatting(chart, dataquestion, ic):
    chart.font.name = 'Arial'
    chart.font.color.theme_color = MSO_THEME_COLOR.TEXT_1
    chart.font.size = Pt(12)
    chart.has_title = True
    chart.chart_title.text_frame.text = dataquestion
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(16)
    chart.chart_title.text_frame.paragraphs[0].font.name = 'Georgia'
    chart.chart_title.text_frame.paragraphs[0].font.bold = False

    if ic == 'PIE':
        for sc, series in enumerate(chart.series):
            for i, point in enumerate(series.points):
                point = chart.series[sc].points[i]
                fill = point.format.fill
                fill.solid()
                fill.fore_color.rgb = brand_colors[i]
    else:
        for i, series in enumerate(chart.series):
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = brand_colors[i]


def toptwobox_formatting(chart, dataquestion):
    stacked_colors = [RGBColor(0, 99, 150), RGBColor(0, 132, 200), RGBColor(0, 0, 0)]
    chart.plots[0].has_data_labels = True
    chart.font.name = 'Arial'
    chart.font.color.theme_color = MSO_THEME_COLOR.TEXT_1
    chart.font.size = Pt(12)
    chart.has_title = True
    chart.chart_title.text_frame.text = dataquestion
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(16)
    chart.chart_title.text_frame.paragraphs[0].font.name = 'Georgia'
    chart.chart_title.text_frame.paragraphs[0].font.bold = False
    chart.plots[0].vary_by_categories = True
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    for i, series in enumerate(chart.series):
        if i < 2:
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = stacked_colors[i]
        else:
            fill = series.format.fill
            fill.background()
    for sequence_color, series in enumerate(chart.series):
        for point in series.points:
            if sequence_color != 2:
                point.data_label.font.color.theme_color = MSO_THEME_COLOR.BACKGROUND_1
            else:
                point.data_label.font.color.theme_color = MSO_THEME_COLOR.TEXT_1
                point.data_label.position = XL_LABEL_POSITION.INSIDE_BASE


def single_series(chart):  # for charts with only one series
    chart.has_legend = False
    chart.plots[0].vary_by_categories = False


def multi_series(chart, ic, splitcolcheck):  # for charts with more than one series
    if splitcolcheck is True:
        chart.has_legend = False
    else:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    chart.plots[0].vary_by_categories = True
    if 'STACKED' not in ic:
        chart.plots[0].overlap = -50


def sub_element(parent, tagname, **kwargs):  # necessary for table formatting
    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element


def table_formatting(table, shape, ic):
    table.first_row = False
    # Selecting Table Style via XML. May be a better solution down the road
    tbl = shape._element.graphic.graphicData.tbl
    if ic == 'TABLE':
        style_id = '{8EC20E35-A176-4012-BC5E-935CFFF8708E}'
    else:
        style_id = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'
    tbl[0][-1].text = style_id


def tablecellformat(cell):
    # Removing Borders via XML. May be a better solution down the road.
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for lines in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
        ln = sub_element(tcPr, lines, w = '0', cap = 'flat', cmpd = 'sng', algn = 'ctr')
    cell.text_frame.paragraphs[0].font.size = Pt(12)


def differing_data_labels(chart, ic):  # Formats data labels for stacked charts and pie accordingly
    chart.plots[0].has_data_labels = True
    if ic == 'PIE':
        for series in chart.series:
            for sequence_color, point in enumerate(series.points):
                if sequence_color != 3:
                    point.data_label.font.color.theme_color = MSO_THEME_COLOR.BACKGROUND_1
                else:
                    point.data_label.font.color.theme_color = MSO_THEME_COLOR.TEXT_1
    else:
        for sequence_color, series in enumerate(chart.series):
            for point in series.points:
                if sequence_color != 3:
                    point.data_label.font.color.theme_color = MSO_THEME_COLOR.BACKGROUND_1
                else:
                    point.data_label.font.color.theme_color = MSO_THEME_COLOR.TEXT_1


def set_reverse_categories(axis):  # workaround function that replicates "Categories in Reverse Order" UI option in PPT
    ele = axis._element.xpath(r'c:scaling/c:orientation')[0]
    ele.set("val", "maxMin")


def axis_formatting(chart, ic):  # formats axis for charts that can have them
    value_axis = chart.value_axis
    category_axis = chart.category_axis
    if ic == "LINE":
        value_axis.visible = True
        value_axis.has_major_gridlines = True
        value_axis.major_tick_mark = XL_TICK_MARK.NONE
        value_axis.format.line.color.theme_color = MSO_THEME_COLOR.BACKGROUND_2
        value_axis.major_gridlines.format.line.color.theme_color = MSO_THEME_COLOR.BACKGROUND_2
    else:
        value_axis.visible = False
        value_axis.has_major_gridlines = False
        value_axis.has_minor_gridlines = False
    category_axis.has_major_gridlines = False
    category_axis.major_tick_mark = XL_TICK_MARK.NONE
    category_axis.format.line.color.theme_color = MSO_THEME_COLOR.BACKGROUND_2
    if ic != 'COLUMN':
        set_reverse_categories(category_axis)


def axis_formatting_100p(chart):  # formats axis for stacked bar and column charts
    value_axis = chart.value_axis  # Turns off value axis
    category_axis = chart.category_axis
    value_axis.visible = False
    value_axis.has_major_gridlines = False
    value_axis.has_minor_gridlines = False
    value_axis = chart.value_axis
    value_axis.maximum_scale = 1
    category_axis.has_major_gridlines = False
    category_axis.major_tick_mark = XL_TICK_MARK.NONE
    category_axis.format.line.color.theme_color = MSO_THEME_COLOR.BACKGROUND_2


def addtextlegend(si, series, mc):
    logging.debug('Text Legend Added', str(mc))
    left = Inches(3.58)
    top = Inches(6.4)
    width = Inches(6)
    height = Inches(0.77)
    shape = prs.slides[si + mc].shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill = shape.fill
    fill.background()
    line = shape.line
    line.fill.background()
    shape.shadow.inherit = False
    for colorsequence, serie in enumerate(series):
        p = shape.text_frame.paragraphs[0]
        p.font.name = 'Arial'
        p.font.bold = False
        p.font.size = Pt(12)
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        font = run.font
        run.text = u'\u25A0' + ' '
        font.color.rgb = brand_colors[colorsequence]
        run = p.add_run()
        font = run.font
        run.text = serie + '\n'
        font.color.rgb = RGBColor(0, 0, 0)


def notesinsert(slide, notestext):
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = str(notestext)


def slidecopyinsert(slide, copyidx, copy):
    slide.placeholders[copyidx[0]].text_frame.text = copy


def createfooter(slide, footercopy, directionalcheck):
    footer = None
    pageno = None
    for shape in slide.placeholders:
        if "SLIDE_NUMBER" in str(shape.placeholder_format.type):  # Master must have placeholder to function
            pageno = slide.placeholders[shape.placeholder_format.idx]
        elif "FOOTER" in str(shape.placeholder_format.type):  # Master must have placeholder to function
            footer = slide.placeholders[shape.placeholder_format.idx]

    footer_text_frame = footer.text_frame
    footer_pageno = pageno.text_frame.paragraphs[0]._p
    # edits XML directly—not ideal, but can't be updated until Python-pptx adds support.
    fld_xml = (
        '<a:fld %s id="{1F4E2DE4-8ADA-4D4E-9951-90A1D26586E7}" type="slidenum">\n'
        '  <a:rPr lang="en-US" smtClean="0"/>\n'
        '  <a:t>2</a:t>\n'
        '</a:fld>\n' % nsdecls("a")
    )
    fld = parse_xml(fld_xml)
    footer_pageno.append(fld)
    footer_a = 'Note: Due to small base, data is directional'
    paragraph_strs = footercopy
    if directionalcheck is True:
        paragraph_strs.append(footer_a)

    for para_str in paragraph_strs:
        p = footer_text_frame.add_paragraph()
        p.text = para_str.replace("'", "")


def footer_patch(self):  # defines new parameters for adding latent footer placeholders
    # Generate a reference to each layout placeholder on this slide layout
    # that should be cloned to a slide when the layout is applied to that
    # slide.
    for ph in self.placeholders:
        yield ph


def assignchartdata(forcefloat, percentcheck, forcecurrency, df):
    chart_data = ChartData()
    chart_data.categories = df.index
    for col in df.columns:
        if percentcheck is True:
            chart_data.add_series(col, df[col], '0%')
        elif forcefloat is True:
            chart_data.add_series(col, df[col], '0.0')
        elif forcecurrency is True:
            chart_data.add_series(col, df[col], '$0.00')
        else:
            chart_data.add_series(col, df[col], '0.0')
    return chart_data


def preflightaddshape(si, msg, lvl, boxplacement):
    left = top = Inches(boxplacement)
    width = height = Inches(2.5)
    shape = prs.slides[si].shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill = shape.fill
    fill.solid()
    if lvl == 'high':
        fill.fore_color.rgb = RGBColor(255, 0, 0)
    else:
        fill.fore_color.rgb = RGBColor(255, 255, 0)
    line = shape.line
    line.fill.background()
    shape.shadow.inherit = False
    shape.text = msg
    shape.text_frame.paragraphs[0].font.color.theme_color = MSO_THEME_COLOR.TEXT_1
    shape.text_frame.paragraphs[0].font.color.size = Pt(10)


def preflight(el, si):  # Inserts applicable error messages as boxes
    if len(el) > 0:
        for chart in el:
            if len(chart) == 0:
                logging.info("No Errors Found")
            else:
                boxplacement = 1.0
                for error in chart:
                    if error in errordict:
                        preflightaddshape(si, errordict[error][0], errordict[error][1], boxplacement)
                        logging.warning(errordict[error][0])
                    boxplacement += 0.5


def charttypeselect(df):  # TO UPDATE: Move to Class to Automatically Identify
    logging.info("Chart type not fount. Selecting a chart type...")
    charttype = "BAR"  # Default chart type

    if len(df.columns) < 2:  # Looks at number of series to determine chart types. (Pie will not have more than one)
        for col in df.columns:
            # If sum of series is within rounding distance of 100%, makes it pie
            if 0.99 <= sum(df[col]) <= 1.01 and len(df[col]) <= 7:
                charttype = "PIE"
    else:
        valcheck = []
        lencheck = []
        rowsum = df.sum(axis=1).tolist()
        for val in rowsum:
            valcount = 0.99 <= val <= 1.01
            valcheck.append(valcount)

        if False not in valcheck:
            charttype = "100% STACKED BAR"
        else:
            for cat in df.index:
                lencheck.append(len(cat))
            if max(lencheck) <= 10 and len(df.index) <= 3:
                charttype = 'COLUMN'
            elif max(lencheck) <= 3 and len(df.index) <= 12:
                charttype = 'COLUMN'
    logging.info(charttype + "chosen based on data")
    return charttype


def shortestval(lst):
    shortest_val = lst[0]
    for val in lst:
        if len(val) < len(shortest_val):
            shortest_val = val
    return shortest_val


def longestval(lst):
    longest_val = lst[0]
    for val in lst:
        if len(val) > len(longest_val):
            longest_val = val
    return longest_val


def chartcreation(slide_info, ph, sd, chart_data, page, chart):
    intendedchart = sd[page][chart]['1—Chart Info'][0]
    forcefloat = sd[page][chart]['1—Chart Info'][3]
    chart_title = shortestval(sd[page][chart]['1—Chart Info'][1])
    snl = sd[page][chart]['1—Chart Info'][5]

    df = sd[page][chart]['5—Data Frame']

    splitcolcheck = slide_info[5]
    splitpagecheck = slide_info[6]

    if intendedchart in ['TABLE', 'HEATMAP']:  # Adds tables and heatmaps
        shape = ph.insert_table(rows=len(df.index) + 1, cols=len(df.columns) + 1)
        table = shape.table
        cell = table.cell(0, 0)  # sets cell value to clear borders
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tablecellformat(cell)  # clears borders from cell 0,0 and updates point size

        for serie, col in enumerate(df.columns):  # Adds/formats column titles
            scount = serie + 1
            cell = table.cell(0, scount)
            cell.text = str(col)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            tablecellformat(cell)  # removes borders from series names and updates point size
            cell.text_frame.paragraphs[0].font.bold = True

        dfcols = list(df)
        for row, idx in enumerate(df.index):  # Adds/formats row titles
            rcount = row + 1
            cell = table.cell(rcount, 0)
            cell.text = idx
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tablecellformat(cell)

            for i, c in enumerate(dfcols):  # Adds/formats values
                scount = i + 1
                cell = table.cell(rcount, scount)
                val = df[c][row]
                if not forcefloat:
                    decconvert = (str(int(val * 100))) + '%'  # This is a crude way of removing decimals
                else:
                    decconvert = str(val)
                cell.text = decconvert
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                if intendedchart != 'TABLE':
                    cellnumval = float(val)
                    negbrightness = cellnumval  # To Update: Transparency should tie to max/min chart vals
                    posbrightness = 1.0 - cellnumval
                    if cellnumval > 0.5:  # Changes background and font color to aid visualization
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(144, 99, 205)
                        cell.fill.fore_color.brightness = posbrightness
                        if posbrightness > 0.50:
                            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                        else:
                            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                    else:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(255, 130, 0)
                        cell.fill.fore_color.brightness = negbrightness
                        if negbrightness > 0.30:
                            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                        else:
                            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                tablecellformat(cell)
        table_formatting(table, shape, intendedchart)

    else:
        if intendedchart == 'COLUMN':
            graphic_frame = ph.insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data)
            chart = graphic_frame.chart
            # Chartype specific formatting
            general_formatting(chart, chart_title, intendedchart)
            if len(snl) > 1:
                multi_series(chart, intendedchart, splitcolcheck)
            else:
                single_series(chart)
            axis_formatting(chart, intendedchart)
            chart.plots[0].has_data_labels = True

        elif intendedchart in ['STACKED COLUMN', '100% STACKED COLUMN']:
            graphic_frame = ph.insert_chart(XL_CHART_TYPE.COLUMN_STACKED, chart_data)
            chart = graphic_frame.chart
            # Chartype specific formatting
            general_formatting(chart, chart_title, intendedchart)
            multi_series(chart, intendedchart, splitcolcheck)
            axis_formatting(chart, intendedchart)
            differing_data_labels(chart, intendedchart)

        elif intendedchart == 'PIE':
            graphic_frame = ph.insert_chart(XL_CHART_TYPE.DOUGHNUT, chart_data)
            chart = graphic_frame.chart
            # Chartype specific formatting
            general_formatting(chart, chart_title, intendedchart)
            multi_series(chart, intendedchart, splitcolcheck)
            hole_size = '75'
            temppath = '/c:chartSpace/c:chart/c:plotArea/c:doughnutChart/c:holeSize'
            chart.element.xpath(temppath)[0].set('val', hole_size)
            differing_data_labels(chart, intendedchart)

        elif intendedchart == 'BAR':
            graphic_frame = ph.insert_chart(XL_CHART_TYPE.BAR_CLUSTERED, chart_data)
            chart = graphic_frame.chart
            # Chartype specific formatting
            general_formatting(chart, chart_title, intendedchart)
            if len(snl) > 1:
                multi_series(chart, intendedchart, splitcolcheck)
            else:
                single_series(chart)
            axis_formatting(chart, intendedchart)
            if splitcolcheck is True or splitpagecheck is True:
                chart.has_title = False
            chart.plots[0].has_data_labels = True

        elif intendedchart in ['STACKED BAR', '100% STACKED BAR']:
            graphic_frame = ph.insert_chart(XL_CHART_TYPE.BAR_STACKED, chart_data)
            chart = graphic_frame.chart
            # Chartype specific formatting
            general_formatting(chart, chart_title, intendedchart)
            multi_series(chart, intendedchart, splitcolcheck)
            axis_formatting(chart, intendedchart)
            differing_data_labels(chart, intendedchart)
            if splitcolcheck is True or splitpagecheck is True:
                chart.has_title = False

        elif intendedchart == 'OVERALL STACKED BAR':
            graphic_frame = ph.insert_chart(XL_CHART_TYPE.BAR_STACKED, chart_data)
            chart = graphic_frame.chart
            # Chartype specific formatting
            general_formatting(chart, chart_title, intendedchart)
            multi_series(chart, intendedchart, splitcolcheck)
            axis_formatting(chart, intendedchart)
            differing_data_labels(chart, intendedchart)
            if splitcolcheck is True or splitpagecheck is True:
                chart.has_title = False

        elif intendedchart == 'LINE':
            graphic_frame = ph.insert_chart(XL_CHART_TYPE.LINE, chart_data)
            chart = graphic_frame.chart
            # Chartype specific formatting
            general_formatting(chart, chart_title, intendedchart)
            multi_series(chart, intendedchart, splitcolcheck)
            axis_formatting(chart, intendedchart)
            chart.plots[0].has_data_labels = True

        elif intendedchart == 'TOP2BOX':
            findmax = df['Top 2 Box'].tolist()
            axismax = float(max(findmax)) + 0.5
            graphic_frame = ph.insert_chart(XL_CHART_TYPE.BAR_STACKED, chart_data)
            chart = graphic_frame.chart
            toptwobox_formatting(chart, chart_title)
            axis_formatting(chart, intendedchart)
            value_axis = chart.value_axis
            value_axis.maximum_scale = axismax


def layout_collect():  # This collects information about the template and drops it into a dictionary
    layout_index = 0
    layout_dict = {}  # Not global
    for layout in prs.slide_layouts:
        medchart = []
        widechart = []
        slimchart = []
        medtable = []
        widetable = []
        slideno = []
        footer = []
        title = []
        body = []
        picture = []
        for shape in layout.placeholders:
            x = shape.placeholder_format.idx
            if "CHART" in str(shape.placeholder_format.type):
                if shape.width > Inches(7.0):
                    widechart.append(x)
                elif shape.width < Inches(4.0):
                    slimchart.append(x)
                else:
                    medchart.append(x)
            elif "TABLE" in str(shape.placeholder_format.type):
                if shape.width > Inches(7.0):
                    widetable.append(x)
                else:
                    medtable.append(x)
            elif "SLIDE_NUMBER" in str(shape.placeholder_format.type):
                slideno.append(x)
            elif "FOOTER" in str(shape.placeholder_format.type):
                footer.append(x)
            elif "TITLE" in str(shape.placeholder_format.type):
                title.append(x)
            elif "BODY" in str(shape.placeholder_format.type):
                body.append(x)
            elif "PICTURE" in str(shape.placeholder_format.type):
                picture.append(x)

        layout_dict[layout_index] = {'1Col Chart': slimchart,
                                     '2Col Chart': medchart,
                                     '3Col Chart': widechart,
                                     '2Col Tbl': medtable,
                                     '3Col Tbl': widetable,
                                     'Sl No': slideno,
                                     'Footer': footer,
                                     'Title': title,
                                     'Body': body,
                                     'Picture': picture
                                     }
        layout_index += 1
    return layout_dict


def layout_chooser(chartcount, splitcolcheck, hlayout, layout_dict, splitpagecheck):  # TO UPDATE: Add table support
    for i in layout_dict:
        ld = layout_dict[i]
        if chartcount < 2:  # If one chart per page, the following tree
            if splitcolcheck is True:
                if len(ld['1Col Chart']) == 2 and len(ld['2Col Chart']) == 0 and len(ld['3Col Chart']) == 0:
                    sl = prs.slide_layouts[i]
                    copyidx = ld['Body'][0]
                    idx = layout_dict[i]['1Col Chart'][0]
                    idx2 = layout_dict[i]['1Col Chart'][1]
                    idx3 = idx4 = None
            elif hlayout is True:
                if len(ld['3Col Chart']) == 1 and len(ld['2Col Chart']) == 0:
                    sl = prs.slide_layouts[i]
                    copyidx = ld['Body'][0]
                    idx = layout_dict[i]['3Col Chart'][0]
                    idx2 = idx3 = idx4 = None
            else:
                if len(ld['2Col Chart']) == 1 and len(ld['1Col Chart']) == 0:
                    sl = prs.slide_layouts[i]
                    copyidx = ld['Body'][0]
                    idx = layout_dict[i]['2Col Chart'][0]
                    idx2 = idx3 = idx4 = None
        else:  # if multiple charts per page, the following tree
            if chartcount == 2:
                if len(ld['2Col Chart']) == 2:
                    sl = prs.slide_layouts[i]
                    copyidx = ld['Body'][0]
                    idx = layout_dict[i]['2Col Chart'][0]
                    idx2 = layout_dict[i]['2Col Chart'][1]
                    idx3 = idx4 = None
            elif chartcount == 3:
                if len(ld['1Col Chart']) == 3:
                    sl = prs.slide_layouts[i]
                    copyidx = ld['Body'][0]
                    idx = layout_dict[i]['1Col Chart'][0]
                    idx2 = layout_dict[i]['1Col Chart'][1]
                    idx3 = layout_dict[i]['1Col Chart'][2]
                    idx4 = None
            elif chartcount == 4:
                if len(ld['1Col Chart']) == 4:
                    sl = prs.slide_layouts[i]
                    copyidx = ld['Body'][0]
                    idx = layout_dict[i]['1Col Chart'][0]
                    idx2 = layout_dict[i]['1Col Chart'][1]
                    idx3 = layout_dict[i]['1Col Chart'][2]
                    idx4 = layout_dict[i]['1Col Chart'][3]
    try:
        return  [sl, idx, idx2, idx3, idx4, splitcolcheck, splitpagecheck, copyidx]
    except UnboundLocalError:
        logging.warning("ERROR: Slide layout couldn't be selected due to improper worksheet coding")


def combinepreflight(wb):  # Checks tab colors to ensure proper selections before import
    most_recent_tabcolor = None
    combinelst = []
    combinecount = 1
    for wksht in wb.worksheets:
        tabcolor = wksht.sheet_properties.tabColor
        if tabcolor is not None:
            if tabcolor == most_recent_tabcolor:  # compares tab to most recent tab
                combinecount += 1
            else:
                if combinecount > 1:  # Adds count to preflight list
                    combinelst.append(combinecount)
                combinecount = 1  # Resets to one
            most_recent_tabcolor = tabcolor
        else:
            combinecount = 1  # Resets to one



def data_collect(wb, country):  # This collects all data from Excel and uses pandas for editing
    slide_data = OrderedDict()

    dictcount = 0
    while dictcount <= len(wb.sheetnames):
        slide_data[dictcount] = {}
        dictcount += 1

    slidecount = 0

    combinecount = 1  # TO UPDATE: Remove once classes are created
    most_recent_tabcolor = None
    for sheetcount, wksht in enumerate(wb.worksheets):
        noteslst = []
        logging.info(str(wksht) + ': ' + wksht.sheet_state)
        if wksht.sheet_state == "visible":

            noteslst.append(wksht)

            transposecheck = False  # True if transpose command is active
            human_transposecheck = False
            sortcheck = False  # True if sort command is active
            forcepercent = False  # True if force percent command is active
            forcefloat = False  # True if force float command is active
            forcecurrency = False  # True if force currency command is active
            percentcheck = False
            exceptionlist = ['-', '- ', '`', '` ', '*', '[', ']', 'u']  # exceptions for cell values
            errorlist = []  # clears error list for each chart
            data_dict = OrderedDict()  # data repository

            intendedchart = None
            data_question, data_base = [[], []]
            slidecopy = []
            tabcolor = wksht.sheet_properties.tabColor  # Used for combining charts

            series_name_list = []  # aids in data collection and format choices **********CAN THIS BE REPLACED?

            '''
            Look through tabs for commands and copy
            '''
            for row in wksht:
                for cell in row:
                    chartstring = wksht[cell.coordinate].value

                    if cell.fill.start_color.index == '00000000':
                        if type(chartstring) is str:
                            if "*TRANSPOSE" in chartstring.upper():
                                human_transposecheck = True
                            elif "*SORT" in chartstring.upper():
                                sortcheck = True
                            elif "*FORCE PERCENT" in chartstring.upper():
                                forcepercent = True
                            elif "*FORCE FLOAT" in chartstring.upper():
                                forcefloat = True
                            elif "*FORCE CURRENCY" in chartstring.upper():
                                forcecurrency = True
                            elif cell.font.underline == 'single':
                                slidecopy.append(chartstring)

                    elif cell.fill.start_color.index == 6:  # detects based on index color
                        intendedchart = chartstring.upper().strip()
                        if intendedchart not in charttypelist:
                            errorlist.append('404')
                            logging.warning('404 ERROR')
                    elif cell.fill.start_color.index == 4:  # detects data_question for footer and chart title
                        data_question.append(chartstring)
                    elif cell.fill.start_color.index == 5:  # detects base for footer
                        data_base.append(chartstring)

            # Categories Collection
            data_dict['categories'] = []

            for row in wksht:  # detects categories based on index color
                for cell in row:
                    if cell.fill.start_color.index == 7:
                        if cell.value is None:
                            errorlist.append('106')
                            logging.warning('106 ERROR')
                        else:
                            # cat = str(filter(lambda ch: ch not in [':', "'", '&#xA0;'], wksht[cell.coordinate].value))
                            cat = str(cell.value).replace(':', '')
                            data_dict['categories'].append(cat)

            #  Series Collection
            for s in range(2, wksht.max_column + 1):
                series_list = []
                for c in range(4, wksht.max_row + 1):
                    cell = wksht.cell(row=c, column=s)
                    if cell.fill.start_color.index == 8:  # Collects Series Name
                        if cell.value is not None:
                            series_name = str(cell.value).replace(':', '')
                        else:
                            series_name = "ERROR_PLACEHOLDER"
                        series_name_list.append(series_name)

                    elif cell.fill.start_color.index == 9:  # detects data
                        if '%' in cell.number_format:
                            percentcheck = True

                        if cell.value is None:  # creates error for empty cells that have been highlighted
                            errorlist.append('106')
                            logging.warning('106 ERROR')
                        elif forcepercent is True:
                            if cell.value in exceptionlist:
                                percentvalue = 0
                            else:
                                percentvalue = float(cell.value) / 100
                            percentcheck = True
                            series_list.append(percentvalue)
                        elif type(cell.value) != float:
                            if type(cell.value) is int and cell.value > 1:
                                series_list.append(float(cell.value))
                            if cell.value == 1 and percentcheck is True:
                                series_list.append(1.0)
                            elif cell.value == 0 or cell.value in exceptionlist:
                                series_list.append(0.0)
                            elif '=' in str(cell.value):
                                errorlist.append('107')  # TO UPDATE: Read formulas
                                logging.warning('107 ERROR')
                                series_list.append(0.0)
                            elif forcepercent is False:
                                if cell.value > 1 or cell.value is None or cell.value in exceptionlist:
                                    pass
                        else:
                            series_list.append(cell.value)

                if len(series_list) > 0:
                    if len(data_dict['categories']) != len(series_list):
                        errorlist.append('100')
                    data_dict[series_name] = series_list

            # Handles Vertical Series
            # for some reason, the last value becomes the first. I change it after the fact, but that's hacky.
            templist = []
            for name in series_name_list:
                if name not in templist:
                    templist.append(name)

            if len(series_name_list) > len(templist):
                series_name_list = templist
                namelength = len(series_name_list)

                max_len = 0
                max_key = ""
                for key in data_dict:
                    if key != 'categories':
                        curlen = len(data_dict[key])
                        if curlen > max_len:
                            max_key = key
                            max_len = curlen
                tempseries = data_dict[max_key]

                for namecount, name in enumerate(series_name_list):
                    templist = tempseries[namecount::namelength]
                    data_dict[name] = templist
                if '100' in errorlist:
                    errorlist.remove('100')

                # Here's the hacky way of doing it

                tempdata_dict = OrderedDict()
                count3 = 0
                for key in data_dict:
                    if count3 != 1:
                        tempdata_dict[key] = data_dict[key]
                    count3 += 1
                count3 = 0
                for key in data_dict:
                    if count3 == 1:
                        tempdata_dict[key] = data_dict[key]
                    count3 += 1

                data_dict = tempdata_dict

            # Creates Data Frame
            try:
                df = pd.DataFrame.from_dict(data_dict)
                df.set_index('categories', inplace=True)
                if country is not None:  # If this is a country report, it will filter out countries
                    for col in df.columns:
                        if col not in ['Global Average', country]:
                            del df[col]

            except Exception:
                df = pd.DataFrame.from_dict(data_error)
                errorlist.append('100')
                logging.warning('100 ERROR')
                df.set_index('categories', inplace=True)

            # Base Collection—Allows multiple bases
            data_base_list = []
            basecount = len(data_base)
            if basecount > 0:
                if range(35) in data_base:
                    directionalcheck = True
                else:
                    directionalcheck = False
            else:
                directionalcheck = False
            for nameidx, base in enumerate(data_base):
                if basecount == 1:
                    data_base_list.append('Base: ' + str(base))
                else:
                    if country is None:
                        if len(series_name_list) != len(data_base):  # Fails when len series = len categories
                            data_base_list.append(str(data_dict['categories'][nameidx]) + ' base: ' + str(base))
                        else:
                            data_base_list.append(str(series_name_list[nameidx]) + ' base: ' + str(base))
                    else:
                        if series_name_list[nameidx] == country:
                            data_base_list.append(str(series_name_list[nameidx]) + ' base: ' + str(base))
                        elif series_name_list[nameidx] == 'Global Average':
                            data_base_list.append(str(series_name_list[nameidx]) + ' base: ' + str(base))
                        else:
                            pass

            # Uses point system to check if data should be transposed
            transposescore = 0
            # hundredmaxlst = []
            # if len(series_name_list) > len(data_dict['categories']):
            #     transposescore += 3
            # if len(data_dict['categories']) < 2:
            #     transposescore += 3

            # for cat in data_dict:
            #     if cat != 'categories':
            #         if 0.99 <= sum(data_dict[cat]) <= 1.01:
            #             hundredmaxlst.append(True)
            #         else:
            #             hundredmaxlst.append(False)
            # if len(hundredmaxlst) > 1:
            #     if False not in hundredmaxlst:
            #         transposescore += 5
            if human_transposecheck is True:
                transposescore += 5
            # if intendedchart in ["TABLE", 'HEATMAP']:
            #     transposescore -= 3
            if transposescore > 4:
                transposecheck = True

            #  Transposes Data before inserting back into dictionary using Pandas
            if transposecheck is True:
                df_t = df.transpose()
                noteslst.append("Chart Transposed")
                df = df_t  # replaces data with transposed version
                series_name_list = list(df.columns.values)
                logging.info("Chart Transposed")

            #  Sort data with Pandas. Does not work if category/series mismatch error thrown earlier
            if sortcheck is True and '100' not in errorlist:
                if len(data_dict) > 0:
                    colname = df.columns[0]
                    s_note = "Chart sorted by " + str(colname)
                    noteslst.append(s_note)
                    if intendedchart == 'TOP2BOX':
                        df_s = df.sort_values(by=[colname, df.columns[1]], ascending=[False, False])
                    else:
                        df_s = df.sort_values(by=[colname], ascending=[False])
                    df = df_s

            # If chart is Top2Box, move Top2Box series to end
            if intendedchart == "TOP2BOX":
                tempcol = []
                for col in df.columns:
                    tempcol.append(col)
                if "Top 2 Box" in df.columns:
                    tempcol.append(tempcol.pop(tempcol.index("Top 2 Box")))
                    df = df.reindex(columns=tempcol)

            #  Combine Check looks at tab color to see if new tab combined with previous on slide
            if tabcolor is not None:
                if tabcolor == most_recent_tabcolor:  # compares tab to most recent tab
                    combinecount += 1
                    slidecount -= 1  # Keeps slide count accurate for error reporting
                else:
                    combinecount = 1  # Resets to one
                most_recent_tabcolor = tabcolor
            else:
                combinecount = 1  # Resets to one

            if intendedchart is None:  # This checks to see if a chart type has been selected, if not, it choose one
                intendedchart = charttypeselect(df)

            # Cleanup error list
            preflightlist = set(errorlist)

            infolist = [
                intendedchart,
                data_question,
                data_base_list,
                forcefloat,
                forcepercent,
                series_name_list,
                sortcheck,
                directionalcheck,
                percentcheck,
                noteslst,
                forcecurrency
            ]
            slide_data[slidecount][combinecount] = {
                '1—Chart Info': infolist,
                '2—Chart Data': data_dict,
                '3—Error List': preflightlist,
                '4—Slide Copy': slidecopy,
                '5—Data Frame': df
            }
            slidecount += 1

    '''Scrub blanks'''
    scrubber_sd = slide_data.copy()
    for scrubcounter, dat in enumerate(slide_data):
        if len(scrubber_sd[scrubcounter]) == 0:
            del scrubber_sd[scrubcounter]
    slide_data = scrubber_sd
    return slide_data


def data_import(sd, ld, countrycheck):
    splitcolcheck = False
    splitpagecheck = False
    hlayout = False
    multicount = 0

    # Select Layout and drop in charts
    for slideidx, page in enumerate(sd):
        chartcount = len(sd[page])
        directionallst = []
        footercopy = []
        preflighterrors = []

        logging.info('Slide no.' + str(slideidx) + '--There is/are' + str(chartcount) + ' Charts:')
        if chartcount == 1:
            for chart in sd[page]:
                intendedchart = sd[page][chart]['1—Chart Info'][0].strip()
                data_question = 'Q: ' + longestval(sd[page][chart]['1—Chart Info'][1])
                footercopy.append(str(sd[page][chart]['1—Chart Info'][2])[2:-2])
                footercopy.append(data_question)
                forcefloat = sd[page][chart]['1—Chart Info'][3]
                forcecurrency = sd[page][chart]['1—Chart Info'][10]
                directionalcheck = sd[page][chart]['1—Chart Info'][7]
                percentcheck = sd[page][chart]['1—Chart Info'][8]
                noteslst = sd[page][chart]['1—Chart Info'][9]
                hlayout = False
                preflighterrors.append(sd[page][chart]['3—Error List'])
                slidecopy = sd[page][chart]['4—Slide Copy']
                df = sd[page][chart]['5—Data Frame']
                series_name_list = df.columns

                '''Insert colors into brand colors to correctly color specific categories'''
                colorinsertlist = []  # Use this to remove colors once the chart is inserted
                if intendedchart == 'PIE':
                    for idxcount, idx in enumerate(df.index):
                        if idx in ['Total', 'Overall']:
                            brand_colors.insert(idxcount, RGBColor(0, 0, 0))
                            colorinsertlist.append(idxcount)
                        elif 'Other' in idx:
                            brand_colors.insert(idxcount, RGBColor(140, 140, 140))
                            colorinsertlist.append(idxcount)
                else:
                    for colcount, col in enumerate(df.columns):
                        if len(df.columns) > 2:
                            if col in ['Total', 'Overall']:
                                brand_colors.insert(colcount, RGBColor(0, 0, 0))
                                colorinsertlist.append(colcount)
                            elif 'Other' in col:
                                brand_colors.insert(colcount, RGBColor(140, 140, 140))
                                colorinsertlist.append(colcount)

                if intendedchart in ['TABLE', 'HEATMAP']:
                    serieslencheck = max([int(len(col)) for col in df.columns]) > 15
                    if serieslencheck is False:
                        for i in ld:
                            if len(ld[i]['2Col Tbl']) == 1:
                                sl = prs.slide_layouts[i]
                                idx = ld[i]['2Col Tbl'][0]
                                copyidx = ld[i]['Body'][0]
                                slide_info = [sl, idx, None, None, None, False, False, copyidx]
                    else:
                        for i in ld:
                            if len(ld[i]['3Col Tbl']) == 1:
                                sl = prs.slide_layouts[i]
                                idx = ld[i]['3Col Tbl'][0]
                                copyidx = ld[i]['Body'][0]
                                slide_info = [sl, idx, None, None, None, False, False, copyidx]
                    chart_data = assignchartdata(forcefloat, percentcheck, forcecurrency, df)
                    slide = prs.slides.add_slide(slide_info[0])
                    ph = slide.placeholders[slide_info[1]]
                    if 1 <= len(slidecopy) <= 3:
                        slide.shapes.title.text = slidecopy[0]
                        slidecopyinsert(slide, copyidx, slidecopy[1])
                    chartcreation(slide_info, ph, sd, chart_data, page, chart)
                else:
                    splitpagecheck = False
                    splitcolcheck = False
                    catlencheck = max(len(idx) for idx in df.index) > 15
                    if len(df.index) > 10 and intendedchart == "BAR":
                        if len(df.index) > 10:
                            if len(df.columns) >= 3 or len(df.index) > 19:
                                if catlencheck is True:
                                    splitpagecheck = len(df.index) > 10
                                else:
                                    splitcolcheck = len(df.index) > 10

                    # Define Splitpoint for split charts
                    splitpoint = len(df.index) // 2
                    df_1 = df.iloc[:splitpoint]
                    df_2 = df.iloc[splitpoint:]

                    if splitcolcheck is True:
                        slide_info = layout_chooser(chartcount, splitcolcheck, hlayout, ld, splitpagecheck)
                        chart_data = assignchartdata(forcefloat, percentcheck, forcecurrency, df_1)
                        slide = prs.slides.add_slide(slide_info[0])
                        ph = slide.placeholders[slide_info[1]]
                        if 1 <= len(slidecopy) <= 3:
                            slide.shapes.title.text = slidecopy[0]
                            slidecopyinsert(slide, copyidx, slidecopy[1])
                        chartcreation(slide_info, ph, sd, chart_data, page, chart)

                        chart_data = assignchartdata(forcefloat, percentcheck, forcecurrency, df_2)
                        ph = slide.placeholders[slide_info[2]]
                        chartcreation(slide_info, ph, sd, chart_data, page, chart)
                        addtextlegend(slideidx, series_name_list, multicount)

                    elif splitpagecheck is True:
                        slide_info = layout_chooser(chartcount, splitcolcheck, hlayout, ld, splitpagecheck)
                        chart_data = assignchartdata(forcefloat, percentcheck, forcecurrency, df_1)
                        slide = prs.slides.add_slide(slide_info[0])
                        ph = slide.placeholders[slide_info[1]]
                        chartcreation(slide_info, ph, sd, chart_data, page, chart)
                        createfooter(slide, footercopy, directionalcheck)
                        notesinsert(slide, noteslst)

                        slide = prs.slides.add_slide(slide_info[0])
                        ph = slide.placeholders[slide_info[1]]
                        copyidx = [slide_info[7]]
                        if 1 <= len(slidecopy) <= 3:
                            slide.shapes.title.text = slidecopy[0]
                            slidecopyinsert(slide, copyidx, slidecopy[1])
                        chart_data = assignchartdata(forcefloat, percentcheck, forcecurrency, df_2)
                        chartcreation(slide_info, ph, sd, chart_data, page, chart)
                        slideidx += 1

                    else:
                        if intendedchart in ['COLUMN', 'STACKED COLUMN', '100% STACKED COLUMN']:
                            hlayout = True
                        if len(df.index) < 9 and len(df.columns) < 6:
                            slide_info = layout_chooser(chartcount, splitcolcheck, hlayout, ld, splitpagecheck)
                            chart_data = assignchartdata(forcefloat, percentcheck, forcecurrency, df)
                            slide = prs.slides.add_slide(slide_info[0])
                            ph = slide.placeholders[slide_info[1]]
                            copyidx = [slide_info[7]]
                            if 1 <= len(slidecopy) <= 3:
                                slide.shapes.title.text = slidecopy[0]
                                slidecopyinsert(slide, copyidx, slidecopy[1])
                            chartcreation(slide_info, ph, sd, chart_data, page, chart)
                        else:
                            slide_info = layout_chooser(chartcount, splitcolcheck, hlayout, ld, splitpagecheck)
                            chart_data = assignchartdata(forcefloat, percentcheck, forcecurrency, df)
                            slide = prs.slides.add_slide(slide_info[0])
                            ph = slide.placeholders[slide_info[1]]
                            copyidx = [slide_info[7]]
                            if 1 <= len(slidecopy) <= 3:
                                slide.shapes.title.text = slidecopy[0]
                                slidecopyinsert(slide, copyidx, slidecopy[1])
                            chartcreation(slide_info, ph, sd, chart_data, page, chart)

                for coloridx in colorinsertlist:
                    del brand_colors[coloridx]

        else:
            slide_info = layout_chooser(chartcount, splitcolcheck, hlayout, ld, splitpagecheck)
            slide = prs.slides.add_slide(slide_info[0])
            multinoteslst = []

            for chart in sd[page]:
                intendedchart = sd[page][chart]['1—Chart Info'][0].strip()
                logging.info(intendedchart + ': ' + min(sd[page][chart]['1—Chart Info'][1]))
                data_question = max(sd[page][chart]['1—Chart Info'][1])
                footercopy.append(str(sd[page][chart]['1—Chart Info'][2])[2:-2])
                footercopy.append(data_question)
                forcefloat = sd[page][chart]['1—Chart Info'][3]
                forcecurrency = sd[page][chart]['1—Chart Info'][10]
                preflighterrors.append(sd[page][chart]['3—Error List'])
                directionallst.append(sd[page][chart]['1—Chart Info'][7])
                percentcheck = sd[page][chart]['1—Chart Info'][8]
                slidecopy = sd[page][chart]['4—Slide Copy']
                noteslst = sd[page][chart]['1—Chart Info'][9]
                hlayout = False

                df = sd[page][chart]['5—Data Frame']

                multinoteslst.append(noteslst)
                '''Insert colors into brand colors to correctly color specific categories'''
                colorinsertlist = []
                if intendedchart == 'PIE':
                    for idxcount, idx in enumerate(df.index):
                        if idx in ['Total', 'Overall']:
                            brand_colors.insert(catcount, RGBColor(0, 0, 0))
                            colorinsertlist.append(idxcount)
                        elif 'Other' in idx:
                            brand_colors.insert(idxcount, RGBColor(140, 140, 140))
                            colorinsertlist.append(idxcount)

                else:
                    for colcount, col in enumerate(df.columns):
                        if len(df.columns) > 2:
                            if col in ['Total', 'Overall']:
                                brand_colors.insert(colcount, RGBColor(0, 0, 0))
                                colorinsertlist.append(colcount)
                            elif 'Other' in col:
                                brand_colors.insert(colcount, RGBColor(140, 140, 140))
                                colorinsertlist.append(colcount)

                chart_data = assignchartdata(forcefloat, percentcheck, forcecurrency, df)
                ph = slide.placeholders[slide_info[chart]]
                copyidx = [slide_info[7]]
                logging.info(copyidx)
                logging.info(slidecopy)
                logging.info(len(slidecopy))
                if 1 <= len(slidecopy) <= 3:
                    slide.shapes.title.text = slidecopy[0]

                chartcreation(slide_info, ph, sd, chart_data, page, chart)

                for coloridx in colorinsertlist:
                    del brand_colors[coloridx]

            if countrycheck is True:
                multicount += 1
        directionalcheck = True in directionallst
        createfooter(slide, footercopy, directionalcheck)
        notesinsert(slide, noteslst)
        preflight(preflighterrors, slideidx)


def xlsxselect():
    xslxfile = app.openBox(title="Choose Import Data", dirName=None, fileTypes=[('excel worksheets', '*.xlsx')],
                           parent=None, multiple=False, mode='r')
    app.setEntry('xlsxfile', xslxfile, callFunction=False)


def press():
    type = app.getOptionBox('Report Type')
    fileentry1 = app.getEntry('xlsxfile')
    workbook = load_workbook(filename=fileentry1, data_only=True)

    if type == 'Global Navigator Country Reports':
        dsave = app.directoryBox(title='Where should country reports be saved?', dirName=None, parent=None)
        for pct, c in enumerate(countrylist):
            progpct = str(int(pct/len(countrylist)*100)) + '%'  # Calculates overall percent complete
            global prs  # This is hacky, but it works.
            prs = Presentation(templatename)
            statusmsg = 'Creating ' + c + ' report'
            progmsg = 'Overall: ' + progpct + ' complete'
            app.setStatusbar(statusmsg, field=0)
            app.setStatusbar(progmsg, field=1)
            app.topLevel.update()
            layout_dict = layout_collect()
            slide_data = data_collect(workbook, c)
            data_import(slide_data, layout_dict, True)
            pptxname = c + '.pptx'
            fulld = dsave + "/" + pptxname
            prs.save(fulld)

    else:  # Covers general import files
        country = None
        try:
            app.setStatusbar('Analyzing Template', field=0)
            app.topLevel.update()  # Must call each time to update status
            layout_dict = layout_collect()
            app.setStatusbar('Reading Excel', field=0)
            app.setStatusbarBg("light gray", field=0)
            app.topLevel.update()
            slide_data = data_collect(workbook, country)
            app.setStatusbar('Importing Data', field=0)
            app.setStatusbarBg("light gray", field=1)
            app.topLevel.update()
            data_import(slide_data, layout_dict, False)
            pptxname = fileentry1[:-5]
            app.setStatusbar('Saving', field=0)
            app.setStatusbarBg("light gray", field=2)
            app.topLevel.update()
            trusave = app.saveBox(title='Save Import File', fileName=pptxname, dirName=None, fileExt=".pptx", fileTypes=[('PowerPoint', '*.pptx')])
            prs.save(trusave)
            app.setStatusbar('Import Complete', field=0)
            app.setStatusbarBg("light gray", field=3)
        except:
            logging.warning('ERROR: Data import failed')
            app.setStatusbar('Aborted. See Log.', field=0)
            app.topLevel.update()


def dup():
    fileentry2 = app.getEntry('Original')
    fileentry3 = app.getEntry('Destination')
    newname = fileentry3[:-5] + '_colorcoded.xlsx'

    wb_orig = load_workbook(filename=fileentry2, data_only=True)
    wb_dest = load_workbook(filename=fileentry3)
    if len(wb_orig.worksheets) != len(wb_dest.worksheets):
        logging.warning('ERROR: WORKSHEET COUNT MISMATCH')
    else:
        for (ws_o, ws_d) in zip(wb_orig, wb_dest):
            logging.info(ws_o + ws_d)
            if ws_o.sheet_state == "hidden":
                ws_d.sheet_state = "hidden"
            for (row_o, row_d) in zip(ws_o, ws_d):
                for (cell_o, cell_d) in zip(row_o, row_d):
                    co = cell_o.coordinate
                    cd = cell_d.coordinate
                    cs_o = ws_o[co].value
                    ws_d[cd].fill = copy(ws_o[co].fill)
                    if ws_d[cd].fill.start_color.index in [4, 6, 7]:  # copies value of chart titles/dataquestions
                        ws_d[cd].value = ws_o[co].value
                    if type(cs_o) == str:
                        if cs_o.upper() in ["*TRANSPOSE", "*SORT", "*FORCE PERCENT", "*FORCE FLOAT", "*FORCE CURRENCY"]:
                            logging.info("Command Copied")
                            ws_d[cd].value = ws_o[co].value
                        elif cs_o.upper() in charttypelist:
                            logging.info("Chart Type Copied")
                            ws_d[cd].value = ws_o[co].value
        wb_dest.save(newname)
        logging.info('Duplication Complete')


SlideLayout.iter_cloneable_placeholders = footer_patch  # replaces module code with redefined latent placeholder code


'''
App controls
'''

with gui('File Selection', '800x600') as app:
    app.setTitle('PowerPoint Importer')
    app.setFont(14)
    app.setBg("white")
    app.setPadding([20, 20])
    app.setInPadding([20, 20])
    app.setStretch("both")

    app.addLabel('title', ('PowerPoint Data Import \n' + version), row=0, column=0)
    app.addWebLink('View Instructions', "http://google.com", row=0, column=1)

    app.addStatusbar(fields=4)
    app.setStatusbarBg("white")
    app.setStatusbarFg("black")

    app.startTabbedFrame("TabbedFrame", colspan=7, rowspan=1)
    app.setTabbedFrameBg('TabbedFrame', "white")
    app.setTabbedFrameTabExpand("TabbedFrame", expand=True)

    # Data Import
    app.startTab('PowerPoint Importer')
    app.addLabelOptionBox('Report Type', reporttypes, row=0, column=2)
    app.addEntry('xlsxfile', colspan=2, row=1, column=1)
    app.addButton('Select File', xlsxselect, row=1, column=3)
    app.addButton('Begin', press, row=2, column=1, colspan=3)
    app.setButton('Begin', '     Begin     ')
    app.stopTab()

    # Excel Color Duplicator
    app.startTab('Color Duplicator')
    app.addLabel('ins1', 'Select Original Color-Coded File')
    app.addFileEntry('Original')
    app.addLabel('ins2', 'Select Document for Color Duplication')
    app.addFileEntry('Destination')
    app.addButton('Duplicate', dup)
    app.stopTab()

    # Settings Tab—Complex Issue Coming eventually
    # app.startTab('Settings')
    # app.addLabel('tmpltlbl', 'Select Template')
    # app.addFileEntry('Template', colspan = 2)
    # app.setEntryDefault('Template', 'Templates/DATAIMPORT.pptx')
    # app.stopTab()

    app.stopTabbedFrame()

    app.go()