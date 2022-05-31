from math import ceil
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_VERTICAL_ANCHOR
from pptx.enum.chart import XL_MARKER_STYLE, XL_DATA_LABEL_POSITION
from pptx.enum.dml import MSO_THEME_COLOR_INDEX
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.chart.data import ChartData

from pptx_handlers.text_handler import insert_text, ParagraphInstance
from pptx_handlers.xml_handlers import page_number_xml, add_table_line_graph, remove_cell_borders_with_xml
from utilities.utility_functions import get_key_with_matching_parameters
from utilities.style_variables import get_brand_color


def get_layout_by_function(page, template):
    parameters = {
        'IGNITE_TitlePage': [page.function == 'ignite cover'],
        'TT—Full Text': [page.function == 'text'],
        'TT—Subsection Intro, Main Ideas': [page.function == 'intro'],
        'TT_Primary_Chart': [page.function == 'chart'],
        'TT_Primary_Table': [page.function == 'table'],
        'TT_Primary_Chart_&_Text': [page.function == 'chart and text'],
        'TT_6_Chart_&_Text': [page.function == 'TT_6_Chart_&_Text'],
        'TT_Chart_Parent_Child': [page.function == 'parent and child'],
        'TT_3_chart_2_table': [page.function == 'dtv_1'],
        'TT_Two_Chart_Equal': [page.function == 'dtv_2'],
        'TT_Primary_Table_&_Text': [page.function == 'dtv_3', page.function == 'table and text'],
        'TT_3_Chart_Dashboard_Flipped': [page.function == 'lto scorecard'],
        'TT_End Wrapper_w_Photos': [page.function == 'end cap']
    }
    layout_name = get_key_with_matching_parameters(parameters)
    if layout_name is not None:
        return [x.layout for x in template.layouts if x.name == layout_name][0]


def get_layout_by_placeholders(page, template):
    for chart in page.charts:
        print(chart.chart_title)
    for layout in template.layouts:
        if all([
            layout.chart_count >= page.chart_count,
            layout.table_count >= page.table_count
        ]):
            return layout.layout
    return None


def get_layout(page, template):
    layout = get_layout_by_function(page, template)
    if layout is None:
        layout = get_layout_by_placeholders(page, template)
    return layout


def get_placeholder_from_shapes(slide, placeholder_type, used_placeholders=None):
    used_placeholders = [] if used_placeholders is None else used_placeholders
    for shape in slide.shapes:
        if placeholder_type in shape.name.upper() and shape.name not in used_placeholders:
            return shape


def get_chart_placeholder(slide, used_list, chart_type):
    parameters = {
        'TABLE': [chart_type.upper() == 'TABLE'],
        'CHART': [True],
    }
    placeholder_type = get_key_with_matching_parameters(parameters)
    shape = get_placeholder_from_shapes(slide, placeholder_type, used_list)
    return shape


def add_chart_title(chart, chart_meta):
    chart.has_title = False
    if chart_meta.chart_title is not None and len(chart_meta.chart_title) > 0:
        title_text = chart_meta.chart_title
        print(f"Creating {title_text}")
        chart.has_title = True
        insert_text(
            ParagraphInstance(title_text, bold=True, uppercase=True, font_size=11, alignment='center'),
            chart.chart_title.text_frame,
            one_level=True
        )


def style_font(chart, chart_meta):
    chart.font.name = chart_meta.chart_style.font
    chart.font.color.theme_color = get_brand_color(chart_meta.chart_style.font_color)
    chart.font.size = Pt(chart_meta.chart_style.font_size)


def style_value_axis(chart, chart_meta):
    v_axis = chart.value_axis
    v_axis.visible = chart_meta.chart_style.v_axis_visible
    v_axis.has_major_gridlines = chart_meta.chart_style.v_axis_major_gridlines
    v_axis.has_minor_gridlines = chart_meta.chart_style.v_axis_minor_gridlines
    if v_axis.visible:
        v_axis.format.line.color.theme_color = chart_meta.chart_style.v_axis_line_color
    if chart_meta.chart_style.v_axis_maximum is not None:
        v_axis.maximum_scale = chart_meta.chart_style.v_axis_maximum
    if chart_meta.chart_style.v_axis_minimum is not None:
        v_axis.minimum_scale = chart_meta.chart_style.v_axis_minimum
    if chart_meta.top_box:
        max_val = max(chart_meta.df[f'Top {len(chart_meta.df.columns) - 1} Box'].tolist())
        v_axis.maximum_scale = max_val + 0.1


def style_category_axis(chart, chart_meta):
    c_axis = chart.category_axis
    c_axis.visible = chart_meta.chart_style.c_axis_visible
    c_axis.reverse_order = chart_meta.chart_style.c_axis_reverse
    c_axis.has_major_gridlines = chart_meta.chart_style.c_axis_major_gridlines
    c_axis.has_minor_gridlines = chart_meta.chart_style.c_axis_minor_gridlines
    c_axis.format.line.color.theme_color = chart_meta.chart_style.c_axis_line_color
    c_axis.tick_label_position = chart_meta.chart_style.c_axis_label_position
    c_axis.major_tick_mark = chart_meta.chart_style.c_axis_major_tick_mark


def style_axis(chart, chart_meta):
    for f in [style_value_axis, style_category_axis]:
        try:
            f(chart, chart_meta)
        except ValueError:
            pass


def style_overlap_and_gap(chart, chart_meta):
    chart.plots[0].gap_width = chart_meta.chart_style.gap
    chart.plots[0].overlap = chart_meta.chart_style.overlap


def style_legend(chart, chart_meta):
    chart.has_legend = any([
        len(chart_meta.df.columns) > 1 and not chart_meta.chart_style.value_table_legend,
        chart_meta.intended_chart_type in ['PIE', 'DONUT']
    ])

    if chart.has_legend:
        chart.legend.include_in_layout = False
        chart.legend.position = chart_meta.chart_style.legend_location
    if chart_meta.chart_style.value_table_legend:
        add_table_line_graph(chart)


def get_color_break(chart, chart_meta):
    is_pie = chart_meta.intended_chart_type in ['PIE', 'DONUT']
    number = len(chart_meta.df.columns) if is_pie else len(chart_meta.df.index)
    return int(ceil(number * 0.75))


def assign_label_color(chart, chart_meta):
    color_break = get_color_break(chart, chart_meta)
    for series_idx, series in enumerate(chart.series):
        for point_idx, point in enumerate(series.points):
            point_color = point.data_label.font.color
            if chart_meta.intended_chart_type.upper() in ['PIE', 'DONUT']:
                if point_idx < color_break:
                    point_color.theme_color = get_brand_color('WHITE')
            else:
                if series_idx < color_break:
                    point_color.theme_color = get_brand_color('WHITE')


def fill_series_or_point(series_or_point, color):
    series_or_point.format.fill.solid()
    series_or_point.format.fill.fore_color.theme_color = get_brand_color(color)


def style_highlighted_series(chart, chart_meta):
    emphasis = [str(x).upper() for x in chart_meta.emphasis]
    de_emphasis = [str(x).upper() for x in chart_meta.de_emphasis]
    category_names = chart_meta.df.index.tolist()
    if any([len(chart.series) > 1, chart_meta.highlight]):
        for series in chart.series:
            if series.name.upper() in emphasis:
                fill_series_or_point(series, 'BLACK')
            elif series.name.upper() in de_emphasis:
                fill_series_or_point(series, 'GRAY')
            if chart_meta.highlight and len(chart.series) == 1:
                for category, point in zip(category_names, series.points):
                    if category.upper() in emphasis:
                        fill_series_or_point(point, 'BLACK')
                    elif category.upper() in de_emphasis:
                        fill_series_or_point(point, 'GRAY')


def assign_label_text(chart, chart_meta):
    label_text = chart_meta.label_text
    color_break = get_color_break(chart, chart_meta)
    for series_idx, series in enumerate(chart.series):
        label_text_series = label_text[series_idx]
        for point_idx, point in enumerate(series.points):
            color = get_brand_color('WHITE' if series_idx <= color_break else 'BLACK')
            point.data_label.text_frame.paragraphs[0].text = label_text_series[point_idx]
            point.data_label.text_frame.paragraphs[0].font.color.theme_color = color


def style_line_markers(chart):
    marker_styles = [
        XL_MARKER_STYLE.CIRCLE,
        XL_MARKER_STYLE.DASH,
        XL_MARKER_STYLE.DIAMOND,
        XL_MARKER_STYLE.DOT,
        XL_MARKER_STYLE.PLUS,
        XL_MARKER_STYLE.SQUARE,
        XL_MARKER_STYLE.STAR,
        XL_MARKER_STYLE.TRIANGLE,
        XL_MARKER_STYLE.X
    ]
    for series, marker_style in zip(chart.series, marker_styles):
        series.marker.style = marker_style
        series.marker.size = 7


def style_last_series_for_top_box(chart, chart_meta):
    for idx, series in enumerate(chart.series):
        if idx == len(chart.series) - 1:
            series.format.fill.background()
            for point in series.points:
                point.data_label.font.color.theme_color = get_brand_color('BLACK')
                point.data_label.position = XL_DATA_LABEL_POSITION.INSIDE_BASE


def style_plots(chart, chart_meta):
    chart.chart_style = chart_meta.chart_style.chart_color
    chart.plots[0].has_data_labels = chart_meta.chart_style.has_labels and not chart_meta.chart_style.value_table_legend
    chart.plots[0].vary_by_categories = any([
        len(chart_meta.df.columns) > 1,
        chart_meta.chart_style.vary_category_color
    ])
    style_highlighted_series(chart, chart_meta)
    if len(chart_meta.label_text) > 0:
        assign_label_text(chart, chart_meta)
    if chart_meta.chart_style.differing_data_labels:
        assign_label_color(chart, chart_meta)
    if chart_meta.intended_chart_type == 'LINE':
        style_line_markers(chart)
    if chart_meta.top_box:
        style_last_series_for_top_box(chart, chart_meta)


def style_chart(graphic_frame, chart_meta):
    chart = graphic_frame.chart
    formatting = [
        style_font,
        add_chart_title,
        style_axis,
        style_overlap_and_gap,
        style_plots,
        style_legend
    ]
    for f in formatting:
        f(chart, chart_meta)


def format_table_cell(
        cell,
        color=None,
        alignment=PP_PARAGRAPH_ALIGNMENT.CENTER,
        brightness=0,
        transparent=False
):
    # Removing Borders via XML. May be a better solution down the road.
    remove_cell_borders_with_xml(cell)
    cell.text_frame.paragraphs[0].font.size = Pt(11)
    cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    cell.text_frame.paragraphs[0].alignment = alignment
    if transparent:
        cell.fill.background()
    if color is not None:
        cell.fill.solid()
        cell.fill.fore_color.theme_color = color
        cell.fill.fore_color.brightness = brightness
        if brightness < 0.15:
            cell.text_frame.paragraphs[0].font.color.theme_color = MSO_THEME_COLOR_INDEX.BACKGROUND_1


def add_top_left_cell(table, chart_meta):
    has_title = chart_meta.chart_title != chart_meta.chart_question
    title_text = chart_meta.chart_title if has_title else ''
    merge_destination = len(chart_meta.df.columns) if has_title else 0
    transparent = has_title

    cell = table.cell(0, 0)  # sets cell value to clear borders
    cell.text = title_text.upper()  # If there's a title, will drop in text
    cell.merge(table.cell(0, merge_destination))  # If there's a title, will merge top row
    format_table_cell(cell, transparent=transparent)  # clears borders from cell 0,0 and updates point size


def add_column_titles(table, title_displacement, chart_meta):
    for col_idx, col in enumerate(chart_meta.df.columns):  # Adds/formats column titles
        cell = table.cell((0 + title_displacement), col_idx + 1)  # +1 leaves empty space top left
        cell.text = str(col)
        format_table_cell(cell)  # removes borders from series names and updates point size
        cell.text_frame.paragraphs[0].font.bold = True


def apply_chart_style(shape, chart_meta):
    tbl = shape._element.graphic.graphicData.tbl
    if chart_meta.intended_chart_type == 'TABLE':
        style_id = '{8EC20E35-A176-4012-BC5E-935CFFF8708E}'
    else:
        style_id = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'
    tbl[0][-1].text = style_id


def format_table(shape, table_height, chart_meta):
    """
    914400 is the number by which shape sizes are divided to get the number of inches
    """
    table = shape.table
    table.first_row = False
    table.horz_banding = chart_meta.chart_style.h_banding and not chart_meta.heat_map
    table.vert_banding = chart_meta.chart_style.v_banding and not chart_meta.heat_map
    table_height = 5.95 if table_height > 4 else table_height
    row_height = round((table_height), 2) / len(table.rows)
    col_width = ((shape.width / 914400) / (len(table.columns) + 1))

    for row in table.rows:
        row.height = Inches(row_height)
    for col_idx, col in enumerate(table.columns):
        col.width = Inches(col_width)
        if col_idx == 0:
            col.width = Inches(col_width * 2)
    apply_chart_style(shape, chart_meta)
    shape.height = Inches(table_height)


def get_cell_color(chart_meta, cell_index, highlight):
    parameters = {
        'green': [highlight],
        'orange': [all([chart_meta.heat_map, cell_index < 1])],
        'purple': [chart_meta.heat_map]
    }
    color = get_key_with_matching_parameters(parameters)
    return None if color is None else get_brand_color(color)


def get_cell_shade(chart_meta, cell_index):
    parameters = {
        0: [
            chart_meta.highlight,
            all([chart_meta.heat_map, any([cell_index < 0.51, cell_index > 1.49])])
        ],
        min([1 / cell_index, cell_index / 1]): [chart_meta.heat_map],

        1: [True]
    }
    return get_key_with_matching_parameters(parameters)


def get_cell_font_color(chart_meta, cell_val, col_idx, row_idx, highlight):
    index_pair = chart_meta.indexes[col_idx]
    change_val = 0 if chart_meta.change_from_previous is None else chart_meta.change_from_previous.iat[row_idx, col_idx]
    parameters = {
        'green': [
            chart_meta.growth and cell_val > 0,
            cell_val > index_pair[1],
            change_val > 0
        ],
        'orange': [
            chart_meta.growth and cell_val < 0,
            cell_val < index_pair[0],
            change_val < 0
        ],
        'white': [highlight],
        'black': [True]
    }
    return get_key_with_matching_parameters(parameters)


def add_column_values(table, title_displacement, chart_meta):
    for row_idx, row in enumerate(chart_meta.df.index):
        cell = table.cell((row_idx + 1 + title_displacement), 0)
        cell.text = str(row)
        format_table_cell(cell, alignment=PP_PARAGRAPH_ALIGNMENT.LEFT)
        row_max = max(chart_meta.df.iloc[row_idx].values.tolist())
        for col_idx, col in enumerate(chart_meta.df.columns):
            cell = table.cell((row_idx + 1 + title_displacement), col_idx + 1)
            cell_val = chart_meta.df.iat[row_idx, col_idx]
            cell_index = round(cell_val / chart_meta.df[col].mean(), 2)
            highlight = chart_meta.highlight and cell_val == row_max
            font_color = get_cell_font_color(chart_meta, cell_val, col_idx, row_idx, highlight)
            if '%' in chart_meta.number_format:
                cell_format = ''.join(["{:.", str(chart_meta.decimal_places), "%}"])
                cell_format = cell_format.format(cell_val)
            else:
                cell_format = str(cell_val)
            bold = any([x in str(font_color) for x in ['ACCENT_3', 'ACCENT-4']])
            text_instance = ParagraphInstance(cell_format, bold=bold, font_color=font_color)
            insert_text([text_instance], cell.text_frame, one_level=True)

            format_table_cell(
                cell,
                color=get_cell_color(chart_meta, cell_index, highlight=highlight),
                brightness=get_cell_shade(chart_meta, cell_index)
            )


def insert_table(chart_meta, placeholder):
    title_displacement = 1 if chart_meta.chart_title != chart_meta.chart_question else 0
    table_height = placeholder.height / 914400
    shape = placeholder.insert_table(
        rows=(len(chart_meta.df.index) + 1 + title_displacement),
        cols=len(chart_meta.df.columns) + 1
    )
    table = shape.table

    add_top_left_cell(table, chart_meta)
    add_column_titles(table, title_displacement, chart_meta)
    add_column_values(table, title_displacement, chart_meta)
    format_table(shape, table_height, chart_meta)


def insert_chart(chart_meta, placeholder):
    chart_data = ChartData()
    chart_data.categories = chart_meta.df.index
    for col in chart_meta.df.columns:
        chart_data.add_series(col, chart_meta.df[col], chart_meta.excel_number_format)

    chart_code = chart_meta.chart_style.chart_type_code
    try:
        graphic_frame = placeholder.insert_chart(chart_code, chart_data)
        style_chart(graphic_frame, chart_meta)
    except AttributeError as e:
        chart_meta.notes.append(str(e))


def insert_page_number(slide):
    """
    If no placeholder exists, will skip
    """
    page_number_placeholder = get_placeholder_from_shapes(slide, 'SLIDE NUMBER')
    try:
        page_number_placeholder_paragraph_xml = page_number_placeholder.text_frame.paragraphs[0]._p
        page_number_xml_fields = page_number_xml()
        page_number_placeholder_paragraph_xml.append(page_number_xml_fields)
    except AttributeError:
        pass


def add_page_title(slide, title_copy):
    if title_copy is not None and len(title_copy) > 0:
        try:
            slide.shapes.title.text = title_copy
        except AttributeError:
            pass


def add_page_copy(slide, copy):
    if copy is not None and len(copy) > 0:
        copy = {0: copy} if isinstance(copy, list) else copy
        if isinstance(copy, dict):
            used_placeholders = []
            for key, copy_list in copy.items():
                placeholder = get_placeholder_from_shapes(slide, 'TEXT', used_placeholders)
                used_placeholders.append(placeholder.name)
                text = [ParagraphInstance(x, indent=0 if idx == 0 else 1) for idx, x in enumerate(copy_list)]
                insert_text(text, placeholder.text_frame)


def add_footer(slide, footer_copy):
    if footer_copy is not None and len(footer_copy) > 0:
        footer_placeholder = get_placeholder_from_shapes(slide, 'FOOTER')
        footer_copy = [ParagraphInstance(x, font_size=None) for x in footer_copy]
        insert_text(footer_copy, footer_placeholder.text_frame, one_level=True)
    insert_page_number(slide)


def add_notes(slide, notes_copy):
    text_frame = slide.notes_slide.notes_text_frame
    new_notes_copy = []
    for tab_name, notes in notes_copy.items():
        new_notes_copy.append(ParagraphInstance(tab_name, bold=True))
        for note in notes:
            new_notes_copy.append(ParagraphInstance(note, bold=False))
    insert_text(new_notes_copy, text_frame)


def add_pictures(slide, pictures):
    used_placeholders = []
    for picture in pictures:
        placeholder = get_placeholder_from_shapes(slide, 'PICTURE', used_placeholders)
        used_placeholders.append(placeholder.name)
        placeholder.insert_picture(picture)


def add_preflight_shapes(slide, preflight):
    pass


def add_shapes(slide, shapes):
    for shape_instance in shapes:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(shape_instance.left),
            Inches(shape_instance.top),
            Inches(shape_instance.width),
            Inches(shape_instance.height)
        )
        shape.shadow.inherit = False
        if shape_instance.text is not None:
            insert_text(shape_instance.text, shape.text_frame)
        if shape_instance.fill_color is None:
            shape.fill.background()
        elif shape_instance.fill_color is not None:
            shape.fill.solid()
            shape.fill.fore_color.theme_color = get_brand_color(shape_instance.fill_color)
        if shape_instance.line is None:
            shape.line.fill.background()


def create_report(report_data, template):
    """
    This iterates through each page instance of the report outline  class
    It selects a layout based on the outline parameters
    It then drops in the chart data, and formats the chart and slide
    """

    for page in report_data.pages:
        layout = get_layout(page, template)
        slide = template.prs.slides.add_slide(layout)
        used_list = []
        for chart_meta in page.charts:
            placeholder = get_chart_placeholder(slide, used_list, chart_meta.intended_chart_type)
            used_list.append(placeholder.name)
            parameters = {
                insert_table: [chart_meta.intended_chart_type == 'TABLE'],
                insert_chart: [True]
            }
            get_key_with_matching_parameters(parameters)(chart_meta, placeholder)
        add_page_title(slide, page.title)
        add_page_copy(slide, page.copy)
        add_footer(slide, page.footer)
        add_notes(slide, page.notes)
        add_pictures(slide, page.pictures)
        add_shapes(slide, page.shapes)
        add_preflight_shapes(slide, page)

    return template.prs
