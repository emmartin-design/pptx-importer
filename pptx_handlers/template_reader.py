from pptx import Presentation
from pptx.slide import SlideLayout


def footer_patch(self):
    for ph in self.placeholders:
        yield ph

SlideLayout.iter_cloneable_placeholders = footer_patch


def get_preferred_layouts(report_type):
    """
    All report types must have their preferred layouts here.
    The design team can help source these titles from the Template doc
    Flexible imports need options to cover all reasonable contingencies
    """
    preferred = {
        'General Import': [
            'TT_TitlePage',
            'TT_Primary_Chart',
            'TT_Primary_Table',
            'TT_Two_Chart_Equal',
            'TT_Two_Table_Equal',
            'TT_3_Chart_Equal',
            'TT_6_Chart_Equal',
            'TT_2x2_Chart',
            'TT_8_Chart',
        ],
        'Global Navigator Country Reports': [
            'TT_TitlePage',
            'TT_Primary_Chart',
            'TT_Primary_Table',
            'TT_Two_Chart_Equal',
            'TT_Two_Table_Equal',
            'TT_3_Chart_Equal',
            'TT_6_Chart_Equal',
            'TT_2x2_Chart',
            'TT_8_Chart',
        ],
        'LTO Scorecard Report': [
            'IGNITE_TitlePage',
            'TT—Subsection Intro, Main Ideas',
            'TT—Full Text',
            'TT_Primary_Table',
            'TT_Primary_Chart',
            'TT_6_Chart_&_Text',
            'TT_3_Chart_Dashboard_Flipped',
            'TT_Primary_Table_&_Text',
            'TT_End Wrapper_w_Photos'
        ],
        'Value Scorecard Report': [
            'IGNITE_TitlePage',
            'TT—Subsection Intro, Main Ideas',
            'TT—Full Text',
            'TT_Primary_Table',
            'TT_Primary_Chart',
            'TT_6_Chart_&_Text',
            'TT_3_Chart_Dashboard_Flipped',
            'TT_Primary_Table_&_Text',
            'TT_End Wrapper_w_Photos'
        ],
        'DirecTV Scorecard': [
            'TT_3_chart_2_table',
            'TT_Two_Chart_Equal',
            'TT_Primary_Table_&_Text'
        ],
        'Consumer Quarterly KPIs': [
            'IGNITE_TitlePage',
            'TT—Subsection Intro, Main Ideas',
            'TT_Primary_Chart_&_Text',
            'TT_6_Chart_&_Text',
            'TT_Half_Chart_Half_Text',
            'TT_End Wrapper_w_Photos'
        ],
        'C-Store Consumer KPIs': [
            'IGNITE_TitlePage',
            'TT—Subsection Intro, Main Ideas',
            'TT_Primary_Chart_&_Text',
            'TT_6_Chart_&_Text',
            'TT_End Wrapper_w_Photos'
        ],
        'Subway Scorecard': [
            'TT_Chart_Parent_Child'
        ]
    }
    return preferred.get(report_type)


class PPTXPlaceholder:
    """
    A metaclass for PPTX layout objects, helps handle the placeholders when used later
    """

    def __init__(self, shape):
        self.placeholder = shape
        self.idx = shape.placeholder_format.idx
        self.type = str(shape.placeholder_format.type).split()[0]
        self.width = round(shape.width.inches, 2)
        self.height = round(shape.height.inches, 2)
        self.left = round(shape.left.inches, 2)
        self.top = round(shape.top.inches, 2)


class PPTXLayout:
    """
    A metaclass for PPTX Layout objects
    Assign preferred status for reports here
    """

    def __init__(self, layout, layout_idx):
        self.layout = layout
        self.name = layout.name
        self.idx = layout_idx
        self.placeholders = [PPTXPlaceholder(x) for x in layout.placeholders]

        self.chart_count = self.get_type_count('CHART')
        self.table_count = self.get_type_count('TABLE')
        self.body_count = self.get_type_count('BODY')
        self.title_count = self.get_type_count('TITLE')
        self.picture_count = self.get_type_count('PICTURE')

        self.footer = [x for x in self.placeholders if x.type == 'FOOTER']
        self.page_number = [x for x in self.placeholders if x.type == 'SLIDE_NUMBER']

        self.width_test = [x.width for x in self.placeholders if x.type in ['CHART', 'TABLE', 'PICTURE']]
        self.preferred = False

    def get_type_count(self, shape_type):
        return len([x for x in self.placeholders if x.type == shape_type.upper()])


class PPTXTemplate:
    """
    This class is a meta class that holds information about all layouts in the template
    It helps select preferred layouts, track counts of placeholders, and can hold additional info
    It's also where the Presentation instance is stored, which is used to create the PPTX file later
    """

    def __init__(self, template_dir, report_type):
        self.prs = Presentation(template_dir)
        self.report_type = report_type
        self.layouts = [PPTXLayout(x, x_idx) for x_idx, x in enumerate(self.prs.slide_layouts)]
        self.assign_preferred_layouts()

    def assign_preferred_layouts(self):
        preferred_layouts_by_report = get_preferred_layouts(self.report_type)
        for layout in self.layouts:
            layout.preferred = layout.name in preferred_layouts_by_report

        self.layouts = [layout for layout in self.layouts if layout.preferred]
