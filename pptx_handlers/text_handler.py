from pptx.util import Inches, Pt

from utilities.style_variables import get_brand_color
from utilities.utility_functions import get_key_with_matching_parameters
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT


class ParagraphInstance:
    formatting_tags = ['/b', '/B', '/i', '/h1', '/h2', '/h3', '/h4', '/h5', '/h6', '/h7',
                       '/q', '/tag', '#', '/*GREEN', '/*ORANGE', '/clear']

    def __init__(
            self,
            text,
            font_size=None,
            font_color=None,
            indent=0,
            bold=None,
            italic=None,
            uppercase=False,
            alignment=False,
            runs=None
    ):
        self.text = text.upper() if uppercase else text
        self.font_size = font_size
        self.font_color = self.get_font_color() if font_color is None else get_brand_color(font_color)
        self.font_brightness = -0.5 if any(['/tag' in self.text, 'ACCENT_3' in str(self.font_color)]) else 0
        self.indent = indent
        self.bold = self.get_weight(bold, text)
        self.italic = True if '/i' in text else italic
        self.alignment = self.get_alignment() if alignment is None else self.get_alignment(str(alignment))
        self.runs = [] if runs is None else runs

        if '#' in self.text:
            for run in self.text.split('#'):
                self.runs.append(RunInstance(run))

        try:
            indent_tag = [x for x in ['/h1', '/h2', '/h3', '/h4', '/h5', '/h6', '/h7'] if x in self.text][0]
            self.indent = int(indent_tag[-1]) - 1
            self.text = self.text.replace(indent_tag, '')
        except IndexError:
            pass

        self.scrub_formatting_codes()

    def get_font_color(self):
        color_parameters = {
            get_brand_color('gray'): ['/tag' in self.text],
            get_brand_color('orange'): ['/*ORANGE' in self.text],
            get_brand_color('green'): ['/*GREEN' in self.text],
        }
        return get_key_with_matching_parameters(color_parameters)

    def get_weight(self, user_selection, text):
        parameters = {
            True: [
                '/b' in text,
                self.font_color is not None and 'TEXT_1' not in str(self.font_color)
            ],
            user_selection: [True]
        }
        return get_key_with_matching_parameters(parameters)

    def scrub_formatting_codes(self):
        scrubbed_text = str(self.text)
        for text_format in self.formatting_tags:
            scrubbed_text = scrubbed_text.replace(text_format, '', -1)
        self.text = scrubbed_text

    @staticmethod
    def get_alignment(align='center'):
        alignments = {
            PP_PARAGRAPH_ALIGNMENT.CENTER: [align == 'center'],
            PP_PARAGRAPH_ALIGNMENT.RIGHT: [align == 'right'],
            PP_PARAGRAPH_ALIGNMENT.LEFT: [True]
        }
        return get_key_with_matching_parameters(alignments)


class RunInstance(ParagraphInstance):
    def __init__(
            self,
            text,
            font_size=None,
            font_color=None,
            indent=0,
            bold=None,
            italic=None,
            uppercase=False,
    ):
        super().__init__(text, font_size, font_color, indent, bold, italic, uppercase)
        self.scrub_formatting_codes()


def text_formatting(text_element, text_instance):
    font = text_element.font
    if text_instance.bold is not None:
        font.bold = text_instance.bold
    if text_instance.italic is not None:
        font.italic = text_instance.italic
    if text_instance.font_color is not None:
        font.color.theme_color = text_instance.font_color
        font.color.brightness = text_instance.font_brightness
    if text_instance.font_size is not None:
        font.size = Pt(text_instance.font_size)
    if text_instance.alignment is not None:
        text_element.alignment = text_instance.alignment


def insert_text(text_instances, text_frame, one_level=False):
    text_instances = [text_instances] if not isinstance(text_instances, list) else text_instances
    for p_idx, paragraph_instance in enumerate(text_instances):
        try:
            p = text_frame.paragraphs[p_idx]
        except IndexError:
            p = text_frame.add_paragraph()

        p.level = paragraph_instance.indent
        if len(paragraph_instance.runs) == 0:
            p.text = paragraph_instance.text
            text_formatting(p, paragraph_instance)
        else:
            for run_instance in paragraph_instance.runs:
                run = p.add_run()
                run.text = run_instance.text
                text_formatting(run, run_instance)
