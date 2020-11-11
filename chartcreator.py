from collections import OrderedDict
from pptx.util import Inches, Pt
from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LABEL_POSITION, XL_LEGEND_POSITION, XL_CHART_TYPE
from pptx.enum.chart import XL_TICK_MARK, XL_DATA_LABEL_POSITION, XL_TICK_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import ChartData
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.slide import SlideLayout
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.oxml.xmlchemy import OxmlElement
import pandas as pd
import numpy as np
import logging

#Created in house
import variables as v


def brightness_level(serieslen, point):
    if serieslen > 2:
        bright_lst = np.linspace(-0.5, 0.5, serieslen)
        brightness = bright_lst[point]
    elif serieslen == 2:
        bright_lst = [-0.5, 0]
        brightness = bright_lst[point]
    else:
        brightness = 0
    return brightness


def topbox_formatting(chart, slen):  # Makes the "overal top box" category transparent
    for i, series in enumerate(chart.series):
        fill = series.format.fill
        if i < (int(slen) - 1):
            pass
        else:
            fill.background()
    for sequence_color, series in enumerate(chart.series):
        for point in series.points:
            if sequence_color != (int(slen) - 1):
                point.data_label.font.color.theme_color = MSO_THEME_COLOR.BACKGROUND_1
            else:
                point.data_label.font.color.theme_color = MSO_THEME_COLOR.TEXT_1
                point.data_label.position = XL_LABEL_POSITION.INSIDE_BASE


def seriescolor_formatting(chart, ic, catlen, slen, other=None, overall=None, highlight=None):
    if ic == 'PIE':
        for sc, series in enumerate(chart.series):
            for i, point in enumerate(series.points):
                if int(catlen) > 2:  # Overwrites color with slate or gray if certain terms used
                    point = chart.series[sc].points[i]
                    fill = point.format.fill
                    fill.solid()
                    if i == overall:
                        fill.fore_color.theme_color = MSO_THEME_COLOR.TEXT_1
                    elif i == other:
                        fill.fore_color.theme_color = MSO_THEME_COLOR.BACKGROUND_2
                        fill.fore_color.brightness = 0.15
    else:
        for i, series in enumerate(chart.series):
            if int(slen) > 2:  # Overwrites color with slate or gray if certain terms used
                fill = series.format.fill
                fill.solid()
                if i == overall:
                    fill.fore_color.theme_color = MSO_THEME_COLOR.TEXT_1
                elif i == other:
                    fill.fore_color.theme_color = MSO_THEME_COLOR.BACKGROUND_2
                    fill.fore_color.brightness = 0.15


def SubElement(parent, tagname, **kwargs):
    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element


def add_table_line_graph(chart):
    plotArea = chart._element.chart.plotArea
    SubElement(plotArea, 'c:dTable') # Add data table tag
    dataTable = chart._element.chart.plotArea.dTable
    # Add values to data table tag
    for sub in ['c:showHorzBorder', 'c:showVertBorder', 'c:showOutline', 'c:showKeys']:
        SubElement(dataTable, sub, val='1')


def colorbreak(number):  # Used to determine when data labels flip to black
    if number > 2:
        if number == 3:
            return 2
        else:
            return round(number / 2) + 1
    else:
        return 1


def assign_highlights(lst):  # Looks at series/category names and records their index for highlighting
    overall, other, highlight = None, None, None
    for item in lst:  # Looks at columns to assign proper index for chart highlights
        itemstr = str(item).upper()
        if itemstr in v.overall_list:
            overall = lst.index(item)
        elif itemstr in v.other_list:
            other = lst.index(item)
        elif itemstr in v.highlight_list:
            highlight = lst.index(item)
    return overall, other, highlight


def differing_data_labels(chart, ic, nos, noc, labeltxt):  # Uses colorbreak to assign label font color point by point
    txtcheck = labeltxt != None
    for series_color, series in enumerate(chart.series):
        for point_color, point in enumerate(series.points):
            pointcolor = point.data_label.font.color
            if txtcheck:
                pointcolor = point.data_label.text_frame.paragraphs[0].runs[0].font.color
                point.data_label.text_frame.word_wrap = False  # Not Working!!!!!!!!!
            pointcolor.theme_color = MSO_THEME_COLOR.TEXT_1
            if ic == 'PIE':
                point.data_label.position = XL_LABEL_POSITION.BEST_FIT
                cbreak = colorbreak(noc)
                if point_color < cbreak:
                    pointcolor.theme_color = MSO_THEME_COLOR.BACKGROUND_1
            else:
                cbreak = colorbreak(nos)
                if series_color < cbreak:
                    pointcolor.theme_color = MSO_THEME_COLOR.BACKGROUND_1


def add_label_text(chart, labeltxt):
    for s_idx, series in enumerate(chart.series):
        for pt_idx, point in enumerate(chart.series[s_idx].points):
            frame = point.data_label.text_frame.paragraphs[0]
            frame.text = labeltxt[s_idx][pt_idx]


def set_reverse_categories(axis):  # workaround function that replicates "Categories in Reverse Order" UI option in PPT
    ele = axis._element.xpath(r'c:scaling/c:orientation')[0]
    ele.set("val", "maxMin")


def sub_element(parent, tagname, **kwargs):  # necessary for table formatting
    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element


def table_formatting(table, shape, ic, banding, height=5.35, width=9.12):
    table.first_row = False
    table.horz_banding = 'rows' in banding
    table.vert_banding = 'cols' in banding
    row_height = round(height, 2) / len(table.rows)  # Creates a 5.35 inch table minimum
    col_width = (width / (len(table.columns) + 1))
    first_col_width =  col_width * 2
    for row in table.rows:
        row.height = Inches(row_height)
    for colidx, col in enumerate(table.columns):
        col.width = Inches(col_width)
        if colidx == 0:
            col.width = Inches(first_col_width)
    # Selecting Table Style via XML. May be a better solution down the road
    tbl = shape._element.graphic.graphicData.tbl
    if ic == 'TABLE':
        style_id = '{8EC20E35-A176-4012-BC5E-935CFFF8708E}'
    else:
        style_id = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'
    tbl[0][-1].text = style_id


def tablecellformat(cell, color=None, alignment=PP_ALIGN.CENTER, highlight=False, brightness=0):
    # Removing Borders via XML. May be a better solution down the road.
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for lines in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
        ln = sub_element(tcPr, lines, w = '0', cap = 'flat', cmpd = 'sng', algn = 'ctr')
    cell.text_frame.paragraphs[0].font.size = Pt(11)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.text_frame.paragraphs[0].alignment = alignment
    if highlight == True:
        cell.fill.solid()
        cell.fill.fore_color.theme_color = color
        cell.fill.fore_color.brightness = brightness
        if brightness < 0.15:
            cell.text_frame.paragraphs[0].font.color.theme_color = MSO_THEME_COLOR.BACKGROUND_1


def legendplace(intendedchart):
    if intendedchart in ['COLUMN', 'STACKED BAR', '100% STACKED BAR']:
        return XL_LEGEND_POSITION.BOTTOM
    else:
        return XL_LEGEND_POSITION.RIGHT


def create_table(df, placeholder, config):
    height, width = placeholder.height.inches, placeholder.width.inches
    shape = placeholder.insert_table(rows=len(df.index) + 1, cols=len(df.columns) + 1)
    table = shape.table
    cell = table.cell(0, 0)  # sets cell value to clear borders
    tablecellformat(cell)  # clears borders from cell 0,0 and updates point size

    for colidx, col in enumerate(df.columns):  # Adds/formats column titles
        cell = table.cell(0, colidx + 1)  # +1 leaves empty space top left
        cell.text = str(col)
        tablecellformat(cell)  # removes borders from series names and updates point size
        cell.text_frame.paragraphs[0].font.bold = True
        colavg = df[col].mean()  # Used to calculate the average of the entire column for heatmap indexing

    dfcols = list(df)
    for rowidx, row in enumerate(df.index):  # Adds/formats row titles
        cell = table.cell(rowidx + 1, 0)
        cell.text = str(row)
        tablecellformat(cell, alignment=PP_ALIGN.LEFT)

        rowmax = None  # Default value for non-highlighted charts

        for colidx, col in enumerate(dfcols):  # Adds/formats values
            cell = table.cell(rowidx + 1, colidx + 1)
            cellval = df[col][rowidx]
            cellindex = round(cellval/colavg, 2)  # For report-interal indexes, not specific report-external indexes

            # Set Color/Shade Defaults
            highlight_cell = False
            color = None
            shade = 1
            cell_text_color = MSO_THEME_COLOR.TEXT_1
            cell_text_brightness = 0
            cell_text_bold = False

            if config['*HIGHLIGHT']:
                highlight_cell = cellval == config['max values'][rowidx]
                color = MSO_THEME_COLOR.ACCENT_3
                shade = 0
            elif config['*HEAT MAP']:
                config['banding'] = ''  # Doesn't add banding
                highlight_cell = True
                shadevals = [1 / cellindex, cellindex / 1]  # The lower of these dictates the cell's shade
                color = MSO_THEME_COLOR.ACCENT_6
                shade = min(shadevals)
                if cellindex < 1:
                    color = MSO_THEME_COLOR.ACCENT_4
                    if cellindex < 0.51:
                        shade = 0
                elif cellindex > 1:
                    if cellindex > 1.49:
                        shade = 0
                else:
                    shade = 0
            elif config['*INDEXES'] is not None:
                indexes = config['*INDEXES']
                if type(indexes) == list:
                    try:
                        index_pair = indexes[colidx]
                        if cellval > index_pair[1]:
                            cell_text_color = MSO_THEME_COLOR.ACCENT_3
                            cell_text_brightness = -.25
                            cell_text_bold = True
                        elif cellval < index_pair[0]:
                            cell_text_color = MSO_THEME_COLOR.ACCENT_4
                            cell_text_bold = True
                    except:
                        pass
                else:
                    try:
                        indexval = indexes[col][rowidx]
                        if indexval > 0:
                            cell_text_color = MSO_THEME_COLOR.ACCENT_3
                            cell_text_brightness = -.25
                            cell_text_bold = True
                        elif indexval < 0:
                            cell_text_color = MSO_THEME_COLOR.ACCENT_4
                            cell_text_bold = True
                    except:
                        pass

            elif config['*GROWTH'] == True:
                if cellval > 0:
                    cell_text_color = MSO_THEME_COLOR.ACCENT_3
                    cell_text_brightness = -.25
                    cell_text_bold = True
                elif cellval < 0:
                    cell_text_color = MSO_THEME_COLOR.ACCENT_4
                    cell_text_bold = True
            if type(cellval) != str:
                if config['percent check']:
                    if config['dec places'] > 0:
                        cell.text = (str(round((float(cellval * 100)), config['dec places']))) + '%'
                    else:
                        try:
                            cell.text = (str(int(cellval * 100))) + '%'
                        except ValueError:  # Handles NaN
                            cell.text = '0%'
                else:
                    cell.text = str(cellval)
            else:
                cell.text = str(cellval)
            font = cell.text_frame.paragraphs[0].font
            font.color.theme_color = cell_text_color
            font.color.brightness = cell_text_brightness
            font.bold = cell_text_bold
            tablecellformat(cell, color=color, highlight=highlight_cell, brightness=shade)
    table_formatting(table, shape, config['intended chart'], config['banding'], height=height, width=width)


def create_chart(df, placeholder, chart_data, config):

    intendedchart = config['intended chart']
    legendloc = config['legend loc']

    graphic_frame = placeholder.insert_chart(v.charttypelist[intendedchart], chart_data)
    chart = graphic_frame.chart
    chart.chart_style = v.chartstyles[config['preferred color']]

    #Handle titles
    if config['chart title'] == None:
        chart.has_title = False
    else:
        chart.has_title = True
        chtitle = chart.chart_title.text_frame  # Chart title font
        chtitle.text = config['chart title'].upper()
        # newtitle = v.titlecleaner(title[0])  # TO UPDATE: Continue Researching text edits.

        chtitle.paragraphs[0].font.size = Pt(11)
        chtitle.paragraphs[0].font.bold = False

    font = chart.font  # Main body font
    font.name = 'Arial'
    font.color.theme_color = MSO_THEME_COLOR.TEXT_1
    font.size = Pt(11)

    chart.plots[0].has_data_labels = config['data labels']

    if config['label text'] != None:
        add_label_text(chart, config['label text'])

    try:  # If chart is not a pie chart, the following will work.
        v_axis, c_axis = chart.value_axis, chart.category_axis

        v_axis.has_major_gridlines, v_axis.has_minor_gridlines = False, False
        v_axis.visible = intendedchart == 'LINE'  # If ic is a line, axis will be visible, else not.

        c_axis.has_major_gridlines = False
        c_axis.major_tick_mark = XL_TICK_MARK.NONE
        c_axis.format.line.color.theme_color = MSO_THEME_COLOR.BACKGROUND_2

        # Override defaults based on chart type
        if intendedchart == 'LINE':
            v_axis.format.line.color.theme_color = MSO_THEME_COLOR.BACKGROUND_2
            c_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
            c_axis.format.line.color.theme_color = MSO_THEME_COLOR.TEXT_1
            add_table_line_graph(chart)
        if not config['cat axis']:
            c_axis.visible = False
            chart.plots[0].gap_width = 70
        if 'COLUMN' not in intendedchart and 'LINE' not in intendedchart:
            set_reverse_categories(c_axis)
        if 'STACKED' in intendedchart:
            v_axis.maximum_scale = 1
        if 'STACKED' not in intendedchart and 'TOP BOX' not in intendedchart:
            chart.plots[0].overlap = -50

        number_of_series = len(df.columns)
        number_of_cats = len(df.index)
        seriesnocheck = number_of_series > 1  # Checks if there are multiple series

        overall, other, highlights = assign_highlights(df.columns.tolist())

    except:  # If chart is a pie chart, the following can happen.
        number_of_series = 1
        number_of_cats = len(df.index)
        seriesnocheck = True
        overall, other, highlight = assign_highlights(df.index.tolist())

    chart.has_legend = seriesnocheck
    if seriesnocheck:
        if legendloc == "DEFAULT":
            chart.legend.position = legendplace(intendedchart)
            chart.legend.include_in_layout = False
        elif legendloc == None:
            chart.has_legend = False
        else:
            try:
                chart.legend.position = v.legend_locations[legendloc]
            except:
                print('EXCEPTION')
                chart.legend.position = legendplace(intendedchart)
                logging.warning('Chart Legend Location not found. Default Used.')
            chart.legend.include_in_layout = False
        if 'STACKED' in intendedchart or intendedchart == 'PIE':
            # This breaks categories showing for some reason
            differing_data_labels(chart, intendedchart, number_of_series, number_of_cats, config['label text'])

    chart.plots[0].vary_by_categories = seriesnocheck
    seriescolor_formatting(chart, intendedchart, str(number_of_cats), str(number_of_series), other, overall, highlight=None)

    if config['*TOP BOX']:
        topbox_formatting(chart, str(len(df.columns)))
